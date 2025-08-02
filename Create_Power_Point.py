from pptx import Presentation
from pptx.util import Inches, Pt
import plotly.graph_objects as go
import plotly.express as px
import pandas as pd
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import random
import os

# Function to generate a sample chart using Plotly (same as in Streamlit)
# Create a new PowerPoint presentation
prs = Presentation()
directory = 'graphs'
timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
font_name = 'Calibri'
image_folder = './images'
images = os.listdir(image_folder)


# Set slide dimensions to 16:9 aspect ratio
prs.slide_width = Inches(26.66)
prs.slide_height = Inches(15)

######### Function to calculate uptime ##########################################################################

def calculate_uptime_percentage(hours, total_hours):
    """
    Calculate the uptime percentage based on total hours (100% - calculated percentage).
    Parameters:
    hours (float): The number of hours to calculate the percentage for.
    total_hours (float): The total hours to base the percentage on (default is 3276).
    260 weekdays per year - 8 elekta holidays = 252 days * 13 hrs (8 a - 9 p)
    for contracts 8 to 5 the total number is: 2268
    Returns:
    float: The remaining percentage (100% - calculated percentage).
    """
    if total_hours == 0:  # Prevent division by zero
        raise ValueError("Total hours cannot be zero.")
    calculated_percentage = round(((hours / total_hours) * 100),1)
    remaining_percentage = round((100 - calculated_percentage),1)
    return round(remaining_percentage,1)

def add_custom_textbox(slide, left:Inches, top:Inches, width: Inches, height: Inches, font_name: str, font_size:Pt, font_color: RGBColor, bold: bool, text: str):
     textbox = slide.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(height))
     text_frame = textbox.text_frame
     text_frame.text = text
     text_frame.paragraphs[0].font.name = font_name
     text_frame.paragraphs[0].font.size = Pt(font_size)
     text_frame.paragraphs[0].font.bold = bold
     text_frame.paragraphs[0].font.color.rgb = font_color

# Add a white rectangle with a colored border in the background of the slide
def add_rectangle_background(slide, left, top, width, height, BGcolor, border):

    # Add the rectangle shape
    rectangle = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )

    # Set the fill color to white
    rectangle.fill.solid()
    rectangle.fill.fore_color.rgb = BGcolor  # White fill

    # Send rectangle to the back of all other shapes
    slide.shapes._spTree.remove(rectangle._element)
    slide.shapes._spTree.insert(2, rectangle._element)

# def generate_chart(monthly_data):

#     fig = px.histogram(monthly_data, x='month', y=['IAAT', 'OAAT'],
#                        barmode='group',
#                        color_discrete_sequence=['rgb(43, 101, 125)', 'rgb(54, 164, 179)'],
#                        nbins=12,
#                        labels={'OAAT': 'Outside Agreed Available Time', 'IAAT': 'Inside Agreed Available Time'})
#     fig.update_layout(bargap=0.5, title='Downtime by Month', yaxis_title='Downtime Hours')
    
#     return fig


def generate_chart(monthly_data):
    # Convert month to string (optional: keeps x-axis readable)
    monthly_data['month'] = pd.to_datetime(monthly_data['month']).dt.strftime('%b %Y')

    # Melt the data to long format
    df_long = monthly_data.melt(id_vars='month', value_vars=['IAAT', 'OAAT'],
                                 var_name='Type', value_name='Hours')

    # Create the bar chart
    fig = px.bar(df_long, x='month', y='Hours', color='Type',
                 barmode='group',
                #  nbins=12,
                 color_discrete_map={
                     'IAAT': 'rgb(43, 101, 125)',
                     'OAAT': 'rgb(54, 164, 179)'
                 },
                 labels={'Hours': 'Downtime Hours', 'month': 'Month'})

    fig.update_layout(
        width=1000,
        height=600,
        bargap=0.7,
        bargroupgap=0.005, 
        title='Monthly Downtime Hours (IAAT vs OAAT)',
        legend_title='Downtime Type',
        template='plotly_white'
    )
    return fig



def create_pie_chart(iaat_hours, total_hours):
        """
        Create a pie chart showing uptime percentage and IAAT.

        Parameters:
        iaat_hours (float): The number of IAAT hours (Inside Agreed Available Time).
        total_hours (float): The total available hours to calculate percentage (default is 3276).

        Returns:
        Plotly pie chart showing the Uptime and IAAT percentages.
        """
        if total_hours == 0:
            raise ValueError("Total hours cannot be zero.")

        # Calculate IAAT percentage and Uptime percentage
        uptime_percentage = (total_hours - iaat_hours) / total_hours * 100
        iaat_percentage = 100 - uptime_percentage
        # Pie chart labels and values
        labels = ['Uptime', 'IAAT']
        values = [uptime_percentage, iaat_percentage]

        # Create pie chart
        fig = go.Figure(data=[go.Pie(labels=labels, values=values)])

        # Add title and customize chart appearance
        fig.update_traces(marker=dict(colors=['rgb(43, 101, 124)', 'rgb(54, 164, 179)']), 
                          textinfo='label+percent', 
                          hoverinfo='label+percent')

        fig.update_layout(
            title_text="Uptime vs Downtime IAAT (Inside Agreed Available Time)",
        )
        return fig

# Function to add Streamlit-like content to PowerPoint slides
def add_slide_with_chart_and_text(slide_title, df, locations, total_hours):

    # Create first slide
    slide_layout = prs.slide_layouts[6] 
    # Title and Content layout (not blank)
    slide = prs.slides.add_slide(slide_layout)

    try:
    # Select random image from folder
        random_image = random.choice(images)

        # image path
        image_path = os.path.join(image_folder, random_image)
    
        # Add the image to the slide
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(26.5), Inches(15))
    
    except Exception as e:
        print(f"Error adding image to slide: {e}")
        print(f"Image path: {image_path}")


    # Customize the text
    # Add a text box
    textbox = slide.shapes.add_textbox(Inches(1), Inches(5.47), Inches(11), Inches(3))  # Define position and size of the textbox
    text_frame = textbox.text_frame
    text_frame.word_wrap = True  # Enable word wrapping
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f'Downtime report last 12 months'
    run.font.size = Pt(90)  # Set font size
    run.font.bold = True  # Make the text bold
    run.font.name = font_name
    run.font.color.rgb = RGBColor(43, 101, 125)  # Set text color to 'rgb(43, 101, 125)'

    # Add title

    # Loope thorugh each location and create slides

    for i in range(1,len(locations), 1):
        current_locations = locations[i:i + 1]
        for j, location in enumerate(current_locations):
            # Filter data based on the current location
                filtered_df = df[df['location'] == location].copy()
                filtered_df['start date'] = pd.to_datetime(filtered_df['start date'], errors='coerce')
                filtered_df['end date'] = pd.to_datetime(filtered_df['end date'], errors='coerce')
            # Group data by month
                filtered_df['month'] = filtered_df['start date'].dt.to_period('M')
                iaat_downtime = round(filtered_df['IAAT'].sum(), 2)
                oaat_downtime = round(filtered_df['OAAT'].sum(), 2)
                
            # Aggregate IAAT and OAAT downtime by month
                monthly_data = filtered_df.groupby('month')[['IAAT', 'OAAT']].sum().reset_index()
                monthly_data['month'] = monthly_data['month'].dt.to_timestamp()
            
            # Make sure graphs images directory exisits
                if not os.path.exists(directory):
                     os.makedirs(directory)

            # Generate histogram
                chart_fig = generate_chart(monthly_data)
                chart_path = os.path.join(directory, f'histogram_downtime_{location}_{timestamp}.png')
                chart_fig.write_image(chart_path)  # Export the chart to an image

            # Generate pie Chart
                pie_fig = create_pie_chart(iaat_downtime,total_hours=total_hours)
                pie_fig_path = os.path.join(directory, f'pie_chart_{location}_{timestamp}.png')
                pie_fig.write_image(pie_fig_path)

            #Create blank slide    
                slide_layout = prs.slide_layouts[6] 
            # Title and Content layout (not blank)
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title 
                add_custom_textbox(slide, 0.25,0.5,24,3,font_name=font_name, font_size=70, font_color=RGBColor(43,101,125), bold=True, text=f'Downtime for {location}')

  
                # Add Total Downtime Metric with rectangle
                # Total Downtime
                # 90Hrs
                # OAAT
                add_custom_textbox(slide, left=5.17,top=3.18,width=24,height=3,font_name=font_name, font_size=25, font_color=RGBColor(43,101,125), bold=True, text=f'Total Downtime')
                add_custom_textbox(slide, left=5.17,top=3.79,width=3.77,height=1.22,font_name=font_name, font_size=80, font_color=RGBColor(43,101,125), bold=True, text=f'{round(iaat_downtime,1)} hrs')
                add_custom_textbox(slide, left=5.17,top=5.16,width=1.69,height=0.37,font_name=font_name, font_size=20, font_color=RGBColor(96,96,96), bold=False, text=f'*OAAT {oaat_downtime} hrs')
                add_rectangle_background(slide,left=Inches(4.85),top=Inches(2.78),width=Inches(4.6),height=Inches(2.91),BGcolor=RGBColor(248,248,248),border=0)

              
                # Add Total Uptime Metric with rectangle
                # Uptime%
                # 99.7% Calculated uptime
                # Target upttime percentage
                add_custom_textbox(slide, left=18.94,top=3.09,width=1.52,height=0.45,font_name=font_name, font_size=25, font_color=RGBColor(43,101,125), bold=True, text=f'Uptime %')
                add_custom_textbox(slide, left=18.94,top=3.66,width=2.91,height=1.22,font_name=font_name, font_size=80, font_color=RGBColor(43,101,125), bold=True, text=f'{calculate_uptime_percentage(iaat_downtime, total_hours)}%')
                add_custom_textbox(slide, left=18.94,top=4.94,width=1.35,height=0.37,font_name=font_name, font_size=20, font_color=RGBColor(96,96,96), bold=False, text=f'Target 97%')
                add_rectangle_background(slide,left=Inches(18.23),top=Inches(2.81),width=Inches(4.6),height=Inches(2.91),BGcolor=RGBColor(248,248,248),border=0)

                
                # Add Charts Histogram and pie chart with a nice withe rectangle background
                chart = slide.shapes.add_picture(chart_path, left=Inches(1.36), top=Inches(6.65), height=Inches(7.3), width=Inches(13.06))
                pie = slide.shapes.add_picture(pie_fig_path, left= Inches(15.87), top=Inches(6.86), height=Inches(6.61),width=Inches(9.47))
                # Add rectangule to slide
                add_rectangle_background(slide,left=Inches(0.22),top=Inches(2.52),width=Inches(26.21),height=Inches(12),BGcolor=RGBColor(255,255,255),border=0.5)
                print(f"Chart image {chart} and Pie {pie} added to slide successfully!")
    
    prs.save(f'{slide_title}_{timestamp}.pptx')


    # Debugging: Check if the image file is saved properly
    try:
        with open(chart_path, 'rb') as f:
            print("Image file saved successfully!")
            
    except FileNotFoundError:
        print("Error: Image file not found!")

    # Add the chart image to the slide at appropriate coordinates
    try:
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(1))
        textbox.text = chart_path
        pic = slide.shapes.add_picture(chart_path, left=Inches(0.3), top=Inches(4.5), height=Inches(2.94), width=Inches(5.26))
        print("Chart image added to slide successfully!")
    except Exception as e:
        print(f"Error adding image to slide: {e}")