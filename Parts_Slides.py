from pptx import Presentation
from pptx.util import Inches, Pt
import plotly.graph_objects as go
import plotly.express as px
import pandas as pd
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import random
import os
from Create_Power_Point import add_rectangle_background


prs = Presentation()
directory = 'graphs/parts'
timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
font_name = 'Calibri'
image_folder = './images'
images = os.listdir(image_folder)
ELEKTA_FONT_COLOR = RGBColor(43,101,125)
RGB_CUSTOM_COLORS = custom_colors = [
    'rgb(43,101,125)',  # Base color (teal-blue)
    'rgb(85,130,145)',  # Lighter and more saturated
    'rgb(25,85,105)',   # Darker and less saturated
    'rgb(135,175,195)', # Much lighter
    'rgb(0,60,80)',     # Very dark teal
    'rgb(58,121,150)',  # Brighter, with more blue
    'rgb(90,140,160)',  # Muted teal-blue
    'rgb(10,70,90)',    # Dark, almost navy
    'rgb(100,180,200)', # Light cyan-blue
    'rgb(30,90,110)'    # Deep teal
]

# Set slide dimensions to 16:9 aspect ratio
prs.slide_width = Inches(26.66)
prs.slide_height = Inches(15)

def add_custom_textbox(slide, left:Inches, top:Inches, width: Inches, height: Inches, font_name: str, font_size:Pt, font_color: RGBColor, bold: bool, text: str):
     textbox = slide.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(height))
     text_frame = textbox.text_frame
     text_frame.text = text
     text_frame.paragraphs[0].font.name = font_name
     text_frame.paragraphs[0].font.size = Pt(font_size)
     text_frame.paragraphs[0].font.bold = bold
     text_frame.paragraphs[0].font.color.rgb = font_color

def generate_parts_slides(title, df, list, df_all_locations):

    # Create first slide
    slide_layout = prs.slide_layouts[6] 
    # Title and Content layout (not blank)
    slide = prs.slides.add_slide(slide_layout)

    try:
    # Select random image from folder
        random_image = random.choice(images)

        # image path
        image_path = os.path.join(image_folder, random_image)
    
        # Add the image to the slide
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(26.5), Inches(15))
    
    except Exception as e:
        print(f"Error adding image to slide: {e}")
        print(f"Image path: {image_path}")

    # Add title to first slide
    add_custom_textbox(slide=slide,
                       left=1,
                       top=5.47,
                       width=11,
                       height=3,
                       font_name=font_name,
                       font_color=ELEKTA_FONT_COLOR,
                       font_size=90,
                       bold=True,
                       text=title)
    
    #####################################################################################################################
    ########################## All Locations ###########################################################################
    ##################################################################################################################3##

     # Create first slide
    slide_layout = prs.slide_layouts[6] 
    # Title and Content layout (not blank)
    slide = prs.slides.add_slide(slide_layout)
    print(df_all_locations)

    # Display metrics for total cost and total number of parts
    total_cost_ip = df_all_locations['total_cost'].sum()
    total_parts_ip = df_all_locations['total_qty'].sum()
    # st.metric(label="Total Cost", value=f"${total_cost_ip:,.2f}")
    # st.metric(label="Total Number of Parts", value=total_parts_ip)

    # Add title
    add_custom_textbox(slide, 
                       1.27,
                       1.08,
                       24,
                       2.9,
                       font_name=font_name, 
                       font_size=70, 
                       font_color=RGBColor(43,101,125), 
                       bold=True, 
                       text='Part Consumption All Locations')
    # Plotly histogram
    fig = px.histogram(df_all_locations, x='created_date', y='total_cost', color='ip',
                   title='Part Consumption by Installed Product',
                   labels={'created_date': 'Date', 'total_qty': 'Total Quantity'},
                   color_discrete_sequence=RGB_CUSTOM_COLORS,
                   nbins=12, barmode='group')
    
    fig.update_layout(width=1000,
                   height=600,
                   bargap=0.5,
                   bargroupgap=0.005)
               
    # Add Part's Total Metric
    add_custom_textbox(slide, 
                       2.59,
                       5.18,
                       4.63,
                       0.83,
                       font_name=font_name, 
                       font_size=40, 
                       font_color=RGBColor(43,101,125), 
                       bold=False, 
                       text="Parts Total")
        # Add Part's Total Metric
    add_custom_textbox(slide, 
                       2.59,
                       6.17,
                       5.49,
                       1,
                       font_name=font_name, 
                       font_size=70, 
                       font_color=RGBColor(43,101,125), 
                       bold=True, 
                       text=f"${total_cost_ip:,.2f}")
    # Add Part's Total Number
    add_custom_textbox(slide, 
                       2.59,
                       8.85,
                       4.63,
                       0.83,
                       font_name=font_name, 
                       font_size=40, 
                       font_color=RGBColor(43,101,125), 
                       bold=False, 
                       text="Number of parts replaced")
        # Add Part's Total Metric
    add_custom_textbox(slide, 
                       2.59,
                       9.85,
                       5.49,
                       1,
                       font_name=font_name, 
                       font_size=70, 
                       font_color=RGBColor(43,101,125), 
                       bold=True, 
                       text=f"{total_parts_ip:,.0f}")
    
    # Set background color and adjust layout for better visibility
    fig.update_layout(
        plot_bgcolor='white',   # Background of the plot
        paper_bgcolor='white',  # Background of the entire figure
        font_color='black',     # Font color for the text
        coloraxis_showscale=True,  # Show color scale
        
    )
    
    
    # Check if the directory exists, if not, create it
    
    if not os.path.exists(directory):
        os.makedirs(directory)

    histogram_path = os.path.join(directory, f'Histogram_all_parts_{timestamp}.png')
    # Try writing the image and add to slide
    try:
        fig.write_image(histogram_path)
        chart = slide.shapes.add_picture(histogram_path, left=Inches(11.46), top=Inches(3.71), height=Inches(10.02), width=Inches(13.35))
        print(f'Image created and added to slide: {chart}')
    except Exception as e:
        print(f'Error writing image: {e}')

    add_rectangle_background(slide,left=Inches(0.81),
                     top=Inches(3),
                     width=Inches(24.75),
                     height=Inches(11.44),
                     BGcolor=RGBColor(248,248,248),
                     border=0)


    #####################################################################################################################
    ########################## Create a slide per Location ##############################################################
    #####################################################################################################################
    # Loop through df and create the slides
    for ip in list:

        #Create a blank slide layout
        slide_layout = prs.slide_layouts[6]
        # add blank layout and crate slide
        slide = prs.slides.add_slide(slide_layout)

        df_ip =df[df['ip'] == ip].copy()
        
        #Group by loaction and date to calculate total quantity and cost
        df_grouped = df_ip.groupby(['location', 'created_date']).agg(
        total_qty=pd.NamedAgg(column='qty', aggfunc='sum'),
        total_cost=pd.NamedAgg(column='price_per_unit', aggfunc=lambda x: (x * df_ip.loc[x.index, 'qty']).sum()) ).reset_index()

        # Display metrics for total cost and total number of parts
        total_cost_ip = df_grouped['total_cost'].sum()
        total_parts_ip = df_grouped['total_qty'].sum()

        # Create histogram for the current IP
        fig = px.histogram(df_grouped, x='created_date', y='total_cost', color='location',
                           title=f'Part Consumption by month.',
                           color_discrete_sequence=['rgb(43, 101, 125)', 'rgb(54, 164, 179)'],
                           labels={'created_date': 'Date', 'total_qty': 'Total Quantity'},
                           nbins=10, barmode='group')
        
        histogram_path = os.path.join(directory, f'Histogram_parts_{timestamp}.png')

        # Calculate and display the three most expensive items
        df_ip.loc[:,'total_item_cost'] = df_ip['price_per_unit'] * df_ip['qty']
        top_5_items = df_ip[['Item', 'total_item_cost']].sort_values(by='total_item_cost', ascending=False).head(5)

      


         # Add title
        add_custom_textbox(slide, 
                           1.27,
                           1.08,
                           24,
                           2.9,
                           font_name=font_name, 
                           font_size=70, 
                           font_color=RGBColor(43,101,125), 
                           bold=True, 
                           text=ip)

        # High Cost Items Textbox        
        add_custom_textbox(slide, 
                           1.5,
                           5.16,
                           6.6,
                           0.84,
                           font_name=font_name, 
                           font_size=60, 
                           font_color=RGBColor(43,101,125), 
                           bold=True, 
                           text='Higher value parts:')
        # High Cost Items Textbox Parts and price       
        result = ""
        initial_top = 6.23
        i = 1
        for index, row in top_5_items.iterrows():  
            add_custom_textbox(slide, 
                    1.5,
                    initial_top,
                    5.49,
                    3.51,
                    font_name=font_name, 
                    font_size=30, 
                    font_color=RGBColor(99,99,99), 
                    bold=False, 
                    text=f"{i}. {row['Item']}: ${row['total_item_cost']:,.2f}\n")
            i += 1
            initial_top += 0.6
        print(result)
 
       # Metric Total Cost of Parts per Location     
        add_custom_textbox(slide, 
                           14.08,
                           3.43,
                           4.17,
                           2.23,
                           font_name=font_name, 
                           font_size=25, 
                           font_color=RGBColor(43,101,125), 
                           bold=False, 
                           text='Total Cost')
        # Metric Total Cost of Parts per Location     
        add_custom_textbox(slide, 
                           14.08,
                           3.95,
                           4.17,
                           2.23,
                           font_name=font_name, 
                           font_size=60, 
                           font_color=RGBColor(43,101,125), 
                           bold=True, 
                           text=f"${total_cost_ip:,.2f}")
        
       # Metric Total Number of Parts per Location     
        add_custom_textbox(slide, 
                           20.47,
                           3.42,
                           4.17,
                           2.23,
                           font_name=font_name, 
                           font_size=25, 
                           font_color=RGBColor(43,101,125), 
                           bold=False, 
                           text='Total # of parts')
        
        # Metric Total Number of Parts per Location     
        add_custom_textbox(slide, 
                           20.47,
                           3.95,
                           4.17,
                           2.23,
                           font_name=font_name, 
                           font_size=60, 
                           font_color=RGBColor(43,101,125), 
                           bold=True, 
                           text=f"{round(total_parts_ip)}")
        
        add_rectangle_background(slide,left=Inches(0.81),
                                 top=Inches(2.7),
                                 width=Inches(24.75),
                                 height=Inches(11.86),
                                 BGcolor=RGBColor(248,248,248),
                                 border=0)

                
        

        # Check if the directory exists, if not, create it
        if not os.path.exists(directory):
            os.makedirs(directory)

        fig.write_image(histogram_path)
        # Add Charts Histogram and pie chart with a nice withe rectangle background
        chart = slide.shapes.add_picture(histogram_path, left=Inches(13.28), top=Inches(5.58), height=Inches(8.42), width=Inches(11.8))
        print(f"Chart image {chart} added to slide successfully!")  

    prs.save(f'{title}_{timestamp}.pptx')

        # Create histogram

