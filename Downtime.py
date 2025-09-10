# Import necessary libraries
import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from datetime import datetime
import io
import base64
import warnings
import requests
import os
import random
import sys
# --- Page and Style Configuration ---

# Set title and configure Streamlit page layout
st.set_page_config(layout="wide", page_title="Downtime Report Generator")

# Suppress specific openpyxl warning
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")

# --- Constants ---
DOWNTIME_URL = 'https://elekta.lightning.force.com/lightning/r/Report/00O6g000006RPtf/view'
QLIK_URL = 'https://qliksense.elekta.com/sense/app/6b74e786-8876-4f7a-8444-478355cc7b84/sheet/401b848c-b981-4694-abae-b1e347a1dfd8/state/analysis'
HOURS_8_TO_9 = 3276  # 252 days * 13 hrs
HOURS_8_TO_5 = 2268  # 252 days * 9 hrs
FONT_NAME = 'Roboto'
IMAGE_FOLDER = 'images'


# --- Styling ---
st.markdown("""
    <style>
    [data-testid="stMetricValue"] {
        font-size: 45px;
        color: rgb(43, 101, 124);
        justify-content: center;
    }
    [data-testid="stMetricDelta"] svg {
        display: none;
    }
    </style>
""", unsafe_allow_html=True)


# --- Helper Functions ---

@st.cache_data
def calculate_uptime_percentage(hours, total_hours):
    """
    Calculate the uptime percentage.
    """
    if total_hours == 0:
        return 100.0
    downtime_percentage = (hours / total_hours) * 100
    uptime_percentage = 100 - downtime_percentage
    return round(uptime_percentage, 1)

@st.cache_data
def create_pie_chart(iaat_hours, total_hours):
    """
    Create a Plotly pie chart for uptime vs. downtime.
    """
    if total_hours == 0 or iaat_hours >= total_hours:
        uptime_percentage = 0
    else:
        uptime_percentage = (total_hours - iaat_hours) / total_hours * 100

    iaat_percentage = 100 - uptime_percentage
    labels = ['Uptime', 'IAAT']
    values = [uptime_percentage, iaat_percentage]

    fig = go.Figure(data=[go.Pie(labels=labels, values=values)])
    fig.update_traces(
        marker=dict(colors=['rgb(43, 101, 124)', 'rgb(54, 164, 179)']),
        textinfo='label+percent',
        hoverinfo='label+percent'
    )
    fig.update_layout(title_text="Uptime vs Downtime IAAT")
    return fig

@st.cache_data
def create_bar_chart(monthly_data):
    """
    Create a Plotly bar chart for monthly downtime.
    """
    monthly_data['month'] = pd.to_datetime(monthly_data['month']).dt.strftime('%b %Y')
    df_long = monthly_data.melt(id_vars='month', value_vars=['IAAT', 'OAAT'],
                                var_name='Type', value_name='Hours')

    fig = px.bar(df_long, x='month', y='Hours', color='Type',
                 barmode='group',
                 color_discrete_map={'IAAT': 'rgb(43, 101, 125)', 'OAAT': 'rgb(54, 164, 179)'},
                 labels={'Hours': 'Downtime Hours', 'month': 'Month', 'Type': 'Downtime Type'})
    fig.update_layout(
        bargap=0.4,
        title='Monthly Downtime Hours (IAAT vs OAAT)',
        legend_title='Downtime Type',
        template='plotly_white'
    )
    return fig

def add_custom_textbox(slide, text, left, top, width, height, font_size, bold=False, color=RGBColor(0, 0, 0), add_shadow=False):
    """Helper to add a formatted textbox to a PowerPoint slide."""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = FONT_NAME
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    if add_shadow:
        try:
            # This attribute might not exist in older versions of python-pptx
            font.shadow.inherit = True 
        except AttributeError:
            # If the shadow attribute doesn't exist, just continue without it
            pass
    text_frame.word_wrap = True

def add_rectangle_background(slide, left, top, width, height, rgb_color, add_shadow):
    """Adds a styled rectangle to the slide and sends it to the back."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = rgb_color
    rect.line.fill.background() # No border
    # Send the rectangle to the very back of all other shapes
    rect.shadow.inherit = add_shadow  # Add shadow to the rectangle
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(0, rect._element)


def add_slide_background_image(slide, prs):
    """Adds a random background image to the title slide from a local folder."""
    if os.path.exists(IMAGE_FOLDER):
        try:
            # Filter for .png files only, ignoring subdirectories
            png_files = [f for f in os.listdir(IMAGE_FOLDER) if os.path.isfile(os.path.join(IMAGE_FOLDER, f)) and f.lower().endswith('.png')]
            
            if png_files:
                random_image_file = random.choice(png_files)
                image_path = os.path.join(IMAGE_FOLDER, random_image_file)
                
                with open(image_path, 'rb') as image_stream:
                    pic = slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                    # Send the picture to the back
                    slide.shapes._spTree.remove(pic._element)
                    slide.shapes._spTree.insert(0, pic._element)
            else:
                raise FileNotFoundError("No PNG files found in the images folder.")

        except Exception as e:
            st.warning(f"Could not load a random image. Using solid background. Error: {e}")
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 32, 51) # Fallback dark blue
    else:
        st.warning(f"Image folder '{IMAGE_FOLDER}' not found. Using solid background for title slide.")
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 32, 51) # Fallback dark blue

def generate_powerpoint(df, locations, total_hours, agreement_target):
    """
    Generates a PowerPoint presentation from the dataframe and returns it as a byte stream.
    """
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # --- Title Slide ---
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_slide_background_image(slide, prs)
    # Add a title text box on top of the image with a shadow for visibility
    add_custom_textbox(slide, "Downtime Report Last 120 Days", 0.35, 3, 6, 2, 60, bold=True, color=RGBColor(43, 101, 124), add_shadow=True)


    # --- Data Slides per Location ---
    for location in locations:
        if location == 'All' or location == 'N/A':
            continue

        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add a main title for the slide
        add_custom_textbox(slide, f'Downtime for {location}', 0.5, 0.2, 15, 1, 44, bold=True, color=RGBColor(43, 101, 125))

        # Add background rectangles for visual structure
        add_rectangle_background(slide, Inches(0.5), Inches(1.5), Inches(6.75), Inches(2.5), RGBColor(235, 235, 235), add_shadow=False)
        add_rectangle_background(slide, Inches(8.25), Inches(1.5), Inches(6.75), Inches(2.5), RGBColor(235, 235, 235),add_shadow=False)
        add_rectangle_background(slide, Inches(0.25), Inches(1), Inches(15.25), Inches(7.8), RGBColor(255, 255, 255),add_shadow=True)


        filtered_df = df[df['location'] == location].copy()
        filtered_df['start date'] = pd.to_datetime(filtered_df['start date'], errors='coerce')
        
        iaat_downtime = round(filtered_df['IAAT'].sum(), 2)
        oaat_downtime = round(filtered_df['OAAT'].sum(), 2)
        uptime_perc = calculate_uptime_percentage(iaat_downtime, total_hours)

        # --- Metrics ---
        add_custom_textbox(slide, "Total Downtime", 2.5, 1.8, 4, 0.5, 24, bold=True, color=RGBColor(43, 101, 125))
        add_custom_textbox(slide, f"{iaat_downtime:.1f} hrs", 2.5, 2.3, 4, 1, 48, bold=True, color=RGBColor(43, 101, 125))
        add_custom_textbox(slide, f"*OAAT {oaat_downtime:.1f} hrs", 2.5, 3.3, 4, 0.5, 16, color=RGBColor(96, 96, 96))

        add_custom_textbox(slide, "Uptime", 10.6, 1.8, 4, 0.5, 24, bold=True, color=RGBColor(43, 101, 125))
        add_custom_textbox(slide, f"{uptime_perc}%", 10.6, 2.3, 4, 1, 48, bold=True, color=RGBColor(43, 101, 125))
        add_custom_textbox(slide, f"Target {agreement_target}", 10.6, 3.3, 4, 0.5, 16, color=RGBColor(96, 96, 96))

        # --- Charts ---
        if filtered_df['start date'].notna().any():
            filtered_df['month'] = filtered_df['start date'].dt.to_period('M')
            monthly_data = filtered_df.groupby('month')[['IAAT', 'OAAT']].sum().reset_index()
            monthly_data['month'] = monthly_data['month'].dt.to_timestamp()
            
            bar_fig = create_bar_chart(monthly_data)
            bar_img_bytes = bar_fig.to_image(format="png", width=800, height=450)
            slide.shapes.add_picture(io.BytesIO(bar_img_bytes), Inches(0.5), Inches(4.5), width=Inches(7))
        
        pie_fig = create_pie_chart(iaat_downtime, total_hours)
        pie_img_bytes = pie_fig.to_image(format="png", width=500, height=400)
        slide.shapes.add_picture(io.BytesIO(pie_img_bytes), Inches(9.0), Inches(4.5), width=Inches(5))


    # Save presentation to a byte stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

def get_ppt_download_link(ppt_stream, filename="Downtime_Report.pptx"):
    """Generates a link to download the PowerPoint presentation."""
    b64 = base64.b64encode(ppt_stream.read()).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{filename}">Click here to download the PowerPoint Presentation</a>'


# --- Main Application ---

# --- Sidebar ---
st.sidebar.title('Settings')
uploaded_file = st.sidebar.file_uploader(label='Load data file:', type=['xlsx', 'xls', 'csv'])
st.sidebar.divider()

if not uploaded_file:
    st.markdown("<h1 style='color: rgb(43, 101, 124);'>Load Downtime Report</h1>", unsafe_allow_html=True)
    st.write('1. Go to Reports on CLM and select the appropriate Downtime Report.')
    st.markdown(f"CLM Downtime Report Link: [Downtime Matrix]({DOWNTIME_URL})")
    st.write('2. Select "Export" > "Details Only" > Format ".xlsx".')
    st.write('3. Upload the file using the sidebar on the left.')
    st.write('4. Review the downtime data and compare it with other sources if needed.')
    st.markdown(f"Qlik Sense Link for comparison: [Qlik Sense Dashboard]({QLIK_URL})")
    st.write('5. Configure settings in the sidebar and click "Generate PowerPoint".')

else:
    # --- Data Loading and Processing ---
    try:
        if uploaded_file.name.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(uploaded_file)
        else:
            df = pd.read_csv(uploaded_file, encoding='utf-8', errors='replace')

        # Clean and rename columns
        columns_to_remove = [
            'Case: Installed Product', 'Case: Customer Resolution Statement',
            'Exclude', 'Exclude Reason', 'Case: Opened Date'
        ]
        # Only drop columns that actually exist in the dataframe
        df.drop(columns=[col for col in columns_to_remove if col in df.columns], inplace=True, errors='ignore')
        
        # --- Defensive Column Handling ---
        column_map = {
            'description': 'Case: Description',
            'start date': 'Date Start of Down Time (Customer Time)',
            'start time': 'Start of Down Time (Customer Time)',
            'end date': 'Date End of Down Time (Customer Time)',
            'end time': 'End of Down Time (Customer Time)',
            'location': 'Case: Location',
            'case': 'Case: Case Number',
            'IAAT': 'Downtime In Agreed Available Time',
            'OAAT': 'Downtime Out Agreed Available Time'
        }

        for new_name, old_name in column_map.items():
            if new_name not in df.columns:
                if old_name in df.columns:
                    df.rename(columns={old_name: new_name}, inplace=True)
                else:
                    st.warning(f"Column '{new_name}' or '{old_name}' not found. Defaulting to a placeholder value.")
                    if 'date' in new_name:
                        df[new_name] = pd.NaT
                    elif new_name in ['IAAT', 'OAAT']:
                        df[new_name] = 0
                    else:
                        df[new_name] = 'N/A'

        # Ensure data types are correct after potential creation/renaming
        df['IAAT'] = pd.to_numeric(df['IAAT'], errors='coerce').fillna(0)
        df['OAAT'] = pd.to_numeric(df['OAAT'], errors='coerce').fillna(0)
        df['case'] = df['case'].astype(str)
        df['location'] = df['location'].astype(str)
        
        st.markdown("<h1 style='color: rgb(43, 101, 124);'>Downtime Report</h1>", unsafe_allow_html=True)

        # --- Sidebar Controls for Data Filtering ---
        locations = ['All'] + sorted(df['location'].unique().tolist())
        selected_location = st.sidebar.selectbox('Select Location:', locations)

        calculate_8_to_5 = st.sidebar.checkbox("Downtime 8am to 5pm")
        total_hours = HOURS_8_TO_5 if calculate_8_to_5 else HOURS_8_TO_9

        agreement_map = {"Silver": "95%", "Gold": "97%", "Platinum": "98%"}
        option = st.sidebar.selectbox('Service agreement type:', list(agreement_map.keys()), index=1)
        selected_service_agreement_uptime = agreement_map[option]

        # --- PowerPoint Generation Button ---
        st.sidebar.divider()
        st.sidebar.title('Create PowerPoint')
        if st.sidebar.button('Generate PowerPoint Presentation'):
            with st.spinner('Generating PowerPoint... Please wait.'):
                locations_to_process = [selected_location] if selected_location != 'All' else df['location'].unique().tolist()
                ppt_stream = generate_powerpoint(df, locations_to_process, total_hours, selected_service_agreement_uptime)
                st.sidebar.success("PowerPoint created successfully!")
                st.sidebar.markdown(get_ppt_download_link(ppt_stream), unsafe_allow_html=True)


        # --- Main Page Display ---
        if selected_location == 'All':
            locations_to_display = sorted(df['location'].unique().tolist())
        else:
            locations_to_display = [selected_location]

        for location in locations_to_display:
            if location == 'N/A':
                continue

            with st.container(border=True):
                filtered_d = df[df['location'] == location].copy()
                filtered_d['start date'] = pd.to_datetime(filtered_d['start date'], errors='coerce')
                
                # Check if there are any valid dates before proceeding
                if filtered_d['start date'].notna().any():
                    filtered_d['month'] = filtered_d['start date'].dt.to_period('M')
                    monthly_data = filtered_d.groupby('month')[['IAAT', 'OAAT']].sum().reset_index()
                    monthly_data['month'] = monthly_data['month'].dt.to_timestamp()
                else:
                    monthly_data = pd.DataFrame(columns=['month', 'IAAT', 'OAAT'])


                iaat_downtime = round(filtered_d['IAAT'].sum(), 2)
                oaat_downtime = round(filtered_d['OAAT'].sum(), 2)


                st.markdown(f"<h4 style='color: rgb(43, 101, 124);'>{location}</h4>", unsafe_allow_html=True)

                metric1, metric2 = st.columns(2)
                with metric1:
                    st.metric('Total Downtime (IAAT)', f'{iaat_downtime} hrs', delta=f'*OAAT {oaat_downtime} hrs', delta_color='off')
                with metric2:
                    uptime_perc = calculate_uptime_percentage(iaat_downtime, total_hours)
                    st.metric('Calculated Uptime %', f'{uptime_perc}%', delta=f'{selected_service_agreement_uptime} target')

                tab1, tab2 = st.tabs(["📈 Charts", "🗃 Raw Data"])

                with tab1:
                    if not monthly_data.empty:
                        col1, col2 = st.columns(2)
                        with col1:
                            st.plotly_chart(create_bar_chart(monthly_data), use_container_width=True)
                        with col2:
                            st.plotly_chart(create_pie_chart(iaat_downtime, total_hours), use_container_width=True)
                        st.caption('*IAAT - Inside Agreed Available Time | *OAAT - Outside Agreed Available Time')
                    else:
                        st.info("No valid date data available to display charts.")

                with tab2:
                    st.dataframe(filtered_d)

    except Exception as e:
        st.error(f"An error occurred while processing the file: {e}")
        st.warning("Please ensure the uploaded file is in the correct format and contains the expected columns.")
