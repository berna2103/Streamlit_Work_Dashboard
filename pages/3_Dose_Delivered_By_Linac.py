import pandas as pd
import streamlit as st
import plotly.express as px
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches

# Constants for column names and colors
ALL_MODES = 'Dose Delivered (All Modes)'
CLINICAL_MODE = 'Clinical Dose Delivered'
ELEKTA_FONT_COLOR = RGBColor(43, 101, 125)

# Custom RGB colors for plots
RGB_CUSTOM_COLORS = [
    'rgb(43,101,125)', 'rgb(135,175,195)', 'rgb(255,87,51)',
    'rgb(255,205,86)', 'rgb(25,85,105)', 'rgb(85,130,145)',
    'rgb(0,60,80)', 'rgb(58,121,150)', 'rgb(10,70,90)', 'rgb(75,192,192)'
]

# Set Pandas option to avoid silent downcasting (for future compatibility)
pd.set_option('future.no_silent_downcasting', True)

# File upload widget on the sidebar to allow multiple Excel files
uploaded_files = st.sidebar.file_uploader("Upload Excel Files", type="xlsx", accept_multiple_files=True)

# Initialize empty DataFrames for Treatments, Terminations, and Beam Data
df_treatments, df_terminations, df = pd.DataFrame(), pd.DataFrame(), pd.DataFrame()

# Function to clean and process Treatments data
def process_treatments(df):
    df['S / N'] = df['S / N'].astype(str).str.replace(',', '')
    return df.groupby('S / N')['# of Treatment Sessions'].mean().reset_index()

# Function to clean and process Terminations data
def process_terminations(df):
    df['S / N'] = df['S / N'].astype(str).str.replace(',', '')
    return df.groupby('S / N')['% Abnormal Termination'].mean().reset_index()

# Function to clean and process Beam Data
def process_beam_data(df):
    df.rename(columns={df.columns[1]: '', df.columns[2]: 'All Linacs ', df.columns[3]: 'All Linacs'}, inplace=True)
    new_columns = ['Energy'] + [f"{df.columns[i]} {df.iloc[0, i]}" for i in range(1, df.shape[1])]
    df.columns = new_columns
    df = df.drop(df.index[:4])
    df = df.dropna(subset=[' Technique'])
    df = df.replace({'-': '0'})
    return df

# Function to create and display charts for a given serial number
def display_charts(df, sn, df_treatments, df_terminations):
    columns = [f'{sn} Dose Delivered (All Modes)', f'{sn}.1 Clinical Dose Delivered']
    df[columns] = df[columns].apply(pd.to_numeric, errors='coerce')
    df[f'{sn} Difference'] = df[columns[0]] - df[columns[1]]
    energy_grouped = df.groupby('Energy').agg({columns[0]: 'sum', columns[1]: 'sum', f'{sn} Difference': 'sum'}).reset_index()
    energy_technique_grouped = df.groupby(['Energy', ' Technique']).agg({columns[0]: 'sum', columns[1]: 'sum', f'{sn} Difference': 'sum'}).reset_index()
    energy_technique_grouped = energy_technique_grouped.sort_values(by=['Energy', ' Technique'])
    df_dynamic = df[df[' Technique'] == 'Dynamic']
    df_static = df[df[' Technique'] == 'Static']

    with st.container():
        st.title(sn)
        col1, col2 = st.columns(2)
        with col1:
            if not df_treatments.empty:
                avg_treatments = df_treatments[df_treatments['S / N'] == sn]['# of Treatment Sessions'].mean()
                st.metric(f'Average Daily Treatments', f"{avg_treatments:.2f}")

            fig_histogram = px.histogram(energy_grouped, x='Energy',
                                         y=[columns[0], columns[1], f'{sn} Difference'],
                                         color_discrete_sequence=RGB_CUSTOM_COLORS,
                                         title='MUs by Energy', barmode='group')
            st.plotly_chart(fig_histogram)
        with col2:
            if not df_terminations.empty:
                abnormal_term = df_terminations[df_terminations['S / N'] == sn]['% Abnormal Termination'].mean() * 100
                st.metric(f'% Beam Terminations', f"{abnormal_term:.2f}%")
            fig_histogram = px.histogram(energy_technique_grouped, x=' Technique',
                                         y=[columns[0], columns[1], f'{sn} Difference'],
                                         color=' Technique', color_discrete_sequence=RGB_CUSTOM_COLORS,
                                         title='MUs by Technique', barmode='group')
            st.plotly_chart(fig_histogram)

# Check if any files have been uploaded
if uploaded_files:
    st.markdown("<h1 style='color: rgb(43, 101, 124);'>Linac Beam Analysis</h1>", unsafe_allow_html=True)
    dataframes = {uploaded_file.name.split(".")[0]: pd.read_excel(uploaded_file) for uploaded_file in uploaded_files}

    # Process Treatments data if the 'Treatments' file has been uploaded
    if 'Treatments' in dataframes:
        df_treatments = process_treatments(dataframes['Treatments'])
        # st.write(df_treatments)

    # Process Terminations data if the 'Terminations' file has been uploaded
    if 'Terminations' in dataframes:
        df_terminations = process_terminations(dataframes['Terminations'])
        # st.write(df_terminations)

    # Process Beam Data if the 'Beam Data' file has been uploaded
    if 'Beam Data' in dataframes:
        df = process_beam_data(dataframes['Beam Data'])
        # st.write(df)

        with st.container():
            col1, col2 = st.columns(2)
            with col1:
                total_mu_all_modes = df['All Linacs  Dose Delivered (All Modes)'].sum()
                st.metric('Total MU Delivered All Modes', f"{total_mu_all_modes:,.0f} MUs")
                fig_all_linacs = px.pie(df, names='Energy', values='All Linacs  Dose Delivered (All Modes)',
                                        color_discrete_sequence=RGB_CUSTOM_COLORS, title='Energy Usage (Locations combined)', hole=0.6)
                st.plotly_chart(fig_all_linacs)

            with col2:
                total_mu_clinical = df['All Linacs Clinical Dose Delivered'].sum()
                st.metric('Total MU Clinical Dose', f"{total_mu_clinical:,.0f} MUs")
                fig_clinical = px.pie(df, names=' Technique', values='All Linacs  Dose Delivered (All Modes)',
                                      color_discrete_sequence=RGB_CUSTOM_COLORS, title='Technique Used (Locations combined)', hole=0.6)
                st.plotly_chart(fig_clinical)

        serial_numbers = {col.split(' ')[0] for col in df.columns if col.split(' ')[0].isdigit()}
        for sn in sorted(serial_numbers):
            display_charts(df, sn, df_treatments, df_terminations)



def create_slide(prs, sn, chart, terminations, treatments, button_clicked):
    if button_clicked:
        # Create a new slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title = slide.shapes.title
        title.text = f'Serial Number: {sn}'

        # Add chart
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(6), Inches(4.5)
        slide.shapes.add_picture(chart, x, y, cx, cy)

        # Add terminations and treatments
        left = Inches(1)
        top = Inches(6)
        width = Inches(6)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f'Terminations: {terminations}\nTreatments: {treatments}'

# Create a new presentation
prs = Presentation()

# Add a button to the Streamlit app
button_clicked = st.button('Create PowerPoint Slide')

# Call the function for each serial number
for sn in sorted(serial_numbers):
    # You'll need to generate the chart, terminations, and treatments data here
    chart = '/Users/bernardojimenez/Documents/Web_Development_Projects/Python_Course_Data_Analysis/StreamLit/graphs/parts/Histogram_all_parts_20240913_194853.png'
    terminations = 'terminations_data'
    treatments = 'treatments_data'
    create_slide(prs, sn, chart, terminations, treatments, button_clicked)

# Save the presentation
prs.save('test.pptx')
