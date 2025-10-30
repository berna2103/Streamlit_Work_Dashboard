# Import necessary libraries
import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime, timedelta, date
import numpy as np
import random # For title slide image

# --- Add these imports for PowerPoint Generation ---
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
import plotly.io as pio
import io
import os
# --------------------------------------------------

# --- NEW: Add these imports for AI Integration ---
import ollama
import json
# --------------------------------------------------

# --- Define Image Folder for Title Slide ---
# Create an 'images' folder next to your script and add background images there.
image_folder = './images'
images = []
if os.path.exists(image_folder) and os.path.isdir(image_folder):
    try:
        images = [f for f in os.listdir(image_folder) if os.path.isfile(os.path.join(image_folder, f))]
    except Exception as e:
        st.sidebar.warning(f"Could not read images folder: {e}")
else:
    st.sidebar.warning(f"'images' folder not found next to the script. Title slide will have no background image.")
# -------------------------------------------


# --- Helper Functions for PowerPoint ---
def add_custom_textbox(slide, left: Inches, top: Inches, width: Inches, height: Inches, font_name: str, font_size: Pt, font_color: RGBColor, bold: bool, text: str, alignment=None):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = text # Set text on paragraph
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = font_color
    if alignment:
        p.alignment = alignment
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    return textbox

# --- NEW: Helper function for KPI Cards (CORRECTED) ---
def add_kpi_card(slide, left_inch: Inches, top_inch: Inches, width_inch: Inches, height_inch: Inches, title_text: str, value_text: str, font_name: str, title_font_size=Pt(16), value_font_size=Pt(32)):
    """
    Adds a formatted KPI card (a shape with two text boxes) to a slide.
    """
    try:
        # Add a shadow shape first for depth
        shadow_left = left_inch + Inches(0.05)
        shadow_top = top_inch + Inches(0.05)
        shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, shadow_left, shadow_top, width_inch, height_inch)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(200, 200, 200) # Light gray shadow
        shadow.line.width = Pt(0) # Make line invisible
        shadow.shadow.inherit = False # Disable default shadow

        # Add the main card shape
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_inch, top_inch, width_inch, height_inch)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255) # White
        
        # --- FIX: Set the line's FILL to solid ---
        card.line.fill.solid()
        card.line.fill.fore_color.rgb = RGBColor(220, 220, 220)
        # --- END FIX ---
        
        card.line.width = Pt(1)

        # Add Title Text Box
        title_box_height = Inches(0.5)
        title_box = slide.shapes.add_textbox(left_inch + Inches(0.2), top_inch + Inches(0.15), width_inch - Inches(0.4), title_box_height)
        title_frame = title_box.text_frame
        
        p_title = title_frame.paragraphs[0]
        p_title.text = title_text # Set text on paragraph
        
        p_title.font.name = font_name
        p_title.font.size = title_font_size
        p_title.font.bold = False
        p_title.font.color.rgb = RGBColor(96, 96, 96) # Gray text
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.NONE
        title_frame.margin_bottom = Inches(0)
        title_frame.margin_top = Inches(0)

        # Add Value Text Box
        value_box_top = top_inch + title_box_height - Inches(0.1) # Position value right under title
        value_box = slide.shapes.add_textbox(left_inch + Inches(0.2), value_box_top, width_inch - Inches(0.4), height_inch - title_box_height - Inches(0.1))
        value_frame = value_box.text_frame
        
        p_value = value_frame.paragraphs[0]
        p_value.text = value_text # Set text on paragraph
        
        p_value.font.name = font_name
        p_value.font.size = value_font_size
        p_value.font.bold = True
        p_value.font.color.rgb = RGBColor(43, 101, 125) # Main color
        value_frame.word_wrap = False
        value_frame.auto_size = MSO_AUTO_SIZE.NONE
        value_frame.margin_bottom = Inches(0)
        value_frame.margin_top = Inches(0)
        value_frame.vertical_anchor = MSO_ANCHOR.TOP # Align text to top

        return card
    except Exception as e:
        print(f"Error adding KPI card: {e}")
        return None
# ----------------------------------------------


# --- Main PowerPoint Generation Function (Updated) ---
def generate_powerpoint_report(
    # Figures
    fig_kpi_trend, fig_cost_split,
    fig_tech, fig_loc,
    fig_activity,
    fig_parts_qty, fig_parts_cost,
    fig_case_trend_total, fig_case_heatmap,
    # Report Details
    report_title, date_range_str,
    # Main KPIs
    kpi_total_tcs, kpi_tcs_label,
    kpi_labor_cost, kpi_parts_cost,
    kpi_labor_label, kpi_parts_label,
    kpi_total_events, kpi_avg_tcs,
    kpi_total_hours, kpi_total_parts,
    # Case KPIs
    kpi_total_cases, kpi_avg_cost_case, kpi_avg_visits_case,
    # NEW TALKING POINT DATA
    kpi_total_discount, summary_top_parts, summary_top_issues # <-- NEW PARAMETERS
    ):
    """
    Generates a polished PowerPoint presentation from the Streamlit app's figures.
    """
    prs = Presentation()
    prs.slide_width = Inches(26.66)
    prs.slide_height = Inches(15)
    font_name = 'Arial' # Standard font

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[6] # Using blank layout
    slide = prs.slides.add_slide(slide_layout)

    if images:
        try:
            random_image = random.choice(images)
            image_path = os.path.join(image_folder, random_image)
            slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            bg_picture = slide.shapes[0]
            spTree = slide.shapes._spTree
            spTree.remove(bg_picture._element)
            spTree.insert(2, bg_picture._element)
        except Exception as e:
            print(f"Error adding title slide background image: {e}")
            add_custom_textbox(slide, Inches(1), Inches(1), Inches(24), Inches(2), font_name, Pt(60), RGBColor(0,0,0), True, "Image Error")

    add_custom_textbox(slide, Inches(1.5), Inches(5), Inches(24), Inches(3), font_name, Pt(80), RGBColor(43, 101, 125), True, report_title)
    add_custom_textbox(slide, Inches(1.5), Inches(8), Inches(24), Inches(2), font_name, Pt(40), RGBColor(96, 96, 96), False, f"Data from: {date_range_str}\nGenerated on: {datetime.now().strftime('%Y-%m-%d')}")

    # --- Slide 2: KPI Dashboard Trends (MODIFIED for Talking Points) ---
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "KPI Dashboard: Value & Trends"
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)

    # --- Add KPI Cards to Slide 2 ---
    card_top_inch = Inches(1.8)
    card_width_inch = Inches(5.5)
    card_height_inch = Inches(1.7)
    card_spacing_inch = Inches(0.8)
    
    add_kpi_card(slide, Inches(1.5), card_top_inch, card_width_inch, card_height_inch, kpi_tcs_label, f"${kpi_total_tcs:,.2f}", font_name)
    add_kpi_card(slide, Inches(1.5) + card_width_inch + card_spacing_inch, card_top_inch, card_width_inch, card_height_inch, "Total Service Visits (WOs)", f"{kpi_total_events:,}", font_name)
    add_kpi_card(slide, Inches(1.5) + 2*(card_width_inch + card_spacing_inch), card_top_inch, card_width_inch, card_height_inch, "Avg Total Cost per Visit", f"${kpi_avg_tcs:,.2f}", font_name)
    add_kpi_card(slide, Inches(1.5) + 3*(card_width_inch + card_spacing_inch), card_top_inch, card_width_inch, card_height_inch, "Total Labor Hours", f"{kpi_total_hours:,.1f} h", font_name)

    # --- Add SINGLE LARGE KPI Trend Chart ---
    charts_top_inch = card_top_inch + card_height_inch + Inches(0.3)
    img_trend_bytes = io.BytesIO(fig_kpi_trend.to_image(format="png", width=1600, height=750, scale=3)) # Reduced height to fit text box
    
    pic_width_inch = Inches(18.5) # Smaller chart to make room for talking points
    pic_height_inch = pic_width_inch * (750 / 1600)
    pic_left = Inches(1.5) # Align chart to the left
    pic_top = charts_top_inch
    
    slide.shapes.add_picture(img_trend_bytes, pic_left, pic_top, width=pic_width_inch, height=pic_height_inch)

    # --- NEW: Talking Points for KPI Slide ---
    bullets_box_left = Inches(1.5) + pic_width_inch + Inches(0.5)
    bullets_box_width = Inches(prs.slide_width - bullets_box_left - Inches(0.5))
    bullets_box_top = charts_top_inch
    
    textbox = slide.shapes.add_textbox(bullets_box_left, bullets_box_top, bullets_box_width, Inches(7.0))
    tf = textbox.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "QBR Key Value Insights"
    p.font.name = font_name
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(43, 101, 125)
    
    tf.add_paragraph().text = f"Total Cost Absorption: Your agreement shielded you from a gross cost of ${kpi_total_tcs + kpi_total_discount:,.2f} for labor and parts."
    tf.paragraphs[-1].font.size = Pt(18)
    tf.paragraphs[-1].font.name = font_name
    tf.paragraphs[-1].level = 1

    # tf.add_paragraph().text = f"Budget Certainty: The agreement successfully absorbed ${kpi_total_discount:,.2f} in waived costs/discounts alone."
    # tf.paragraphs[-1].font.size = Pt(18)
    # tf.paragraphs[-1].font.name = font_name
    # tf.paragraphs[-1].level = 1

    tf.add_paragraph().text = f"Operational Continuity: We successfully executed {kpi_total_events:,} service visits, representing {kpi_total_hours:,.1f} hours of covered labor, ensuring maximum uptime."
    tf.paragraphs[-1].font.size = Pt(18)
    tf.paragraphs[-1].font.name = font_name
    tf.paragraphs[-1].level = 1
    # --- END NEW: Talking Points for KPI Slide ---


    # --- Slide 3: KPI Cost Split Chart ---
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "KPI Dashboard: Cost Split (Labor vs. Parts)"
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)

    # Add Cost Split chart (large)
    img_split_bytes = io.BytesIO(fig_cost_split.to_image(format="png", width=1600, height=900, scale=3))

    pic_width_inch = Inches(24) # Make it large
    pic_height_inch = pic_width_inch * (900 / 1600)
    pic_left = (prs.slide_width - pic_width_inch) / 2
    pic_top = (prs.slide_height - pic_height_inch) / 2 + Inches(0.5) # Center vertically below title
    
    slide.shapes.add_picture(img_split_bytes, pic_left, pic_top, width=pic_width_inch, height=pic_height_inch)


    # --- NEW Slide 4: Performance (Technician) ---
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Technician Performance"
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)

    # Add Technician chart (large)
    img_tech_bytes = io.BytesIO(fig_tech.to_image(format="png", width=1200, height=900, scale=3))

    pic_width_inch = Inches(20) # A bit narrower for bar chart
    pic_height_inch = pic_width_inch * (900 / 1200)
    pic_left = (prs.slide_width - pic_width_inch) / 2
    pic_top = (prs.slide_height - pic_height_inch) / 2 + Inches(0.5) 

    slide.shapes.add_picture(img_tech_bytes, pic_left, pic_top, width=pic_width_inch, height=pic_height_inch)


    # --- NEW Slide 5: Performance (Location) ---
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Location Performance"
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)

    # Add Location chart (large)
    img_loc_bytes = io.BytesIO(fig_loc.to_image(format="png", width=1200, height=900, scale=3))
    
    pic_width_inch = Inches(20) # A bit narrower
    pic_height_inch = pic_width_inch * (900 / 1200)
    pic_left = (prs.slide_width - pic_width_inch) / 2
    pic_top = (prs.slide_height - pic_height_inch) / 2 + Inches(0.5)

    slide.shapes.add_picture(img_loc_bytes, pic_left, pic_top, width=pic_width_inch, height=pic_height_inch)


    # --- Slide 6: Case Analysis - Trend & KPIs (was 4) ---
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Case Analysis: Trend (Corrective Service)"
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)

    # Add KPI Cards to Slide 6
    card_top_inch_case = Inches(1.8)
    card_width_inch_case = Inches(6.5)
    card_height_inch_case = Inches(1.7)
    card_spacing_inch_case = Inches(1.0)
    
    add_kpi_card(slide, Inches(2.5), card_top_inch_case, card_width_inch_case, card_height_inch_case, "Total Cases", f"{kpi_total_cases:,}", font_name)
    add_kpi_card(slide, Inches(2.5) + card_width_inch_case + card_spacing_inch_case, card_top_inch_case, card_width_inch_case, card_height_inch_case, f"Avg. Cost per Case ({kpi_parts_label})", f"${kpi_avg_cost_case:,.2f}", font_name)
    add_kpi_card(slide, Inches(2.5) + 2*(card_width_inch_case + card_spacing_inch_case), card_top_inch_case, card_width_inch_case, card_height_inch_case, "Avg. Visits per Case", f"{kpi_avg_visits_case:,.1f}", font_name)

    # Add the Case Trend chart (fig_case_trend_total)
    charts_top_inch_case = card_top_inch_case + card_height_inch_case + Inches(0.3)
    img_case_width_px = 1600
    img_case_height_px = 800
    img_case_bytes = io.BytesIO(fig_case_trend_total.to_image(format="png", width=img_case_width_px, height=img_case_height_px, scale=3))
    
    pic_width_inch = Inches(24) 
    pic_height_inch = pic_width_inch * (img_case_height_px / img_case_width_px)
    if pic_height_inch > Inches(10.5): # Constrain height
        pic_height_inch = Inches(10.5)
        pic_width_inch = pic_height_inch * (img_case_width_px / img_case_height_px)

    pic_left = (prs.slide_width - pic_width_inch) / 2
    pic_top = charts_top_inch_case # Position below cards
    slide.shapes.add_picture(img_case_bytes, pic_left, pic_top, width=pic_width_inch, height=pic_height_inch)


    # --- Slide 7: Case Analysis - Heatmap (was 5) ---
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Case Analysis: New Cases by Location (Heatmap)"
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
    
    # Add the Case Heatmap (fig_case_heatmap)
    img_case_width_px = 1800
    img_case_height_px = 1000
    img_case_bytes = io.BytesIO(fig_case_heatmap.to_image(format="png", width=img_case_width_px, height=img_case_height_px, scale=3))

    pic_width_inch = Inches(24) 
    pic_height_inch = pic_width_inch * (img_case_height_px / img_case_width_px)
    if pic_height_inch > Inches(12): # Don't let it exceed slide height bounds
        pic_height_inch = Inches(12)
        pic_width_inch = pic_height_inch * (img_case_width_px / img_case_height_px)

    pic_left = (prs.slide_width - pic_width_inch) / 2
    pic_top = Inches(2.0) + ((Inches(15-2.0)) - pic_height_inch) / 2 # Center below title
    slide.shapes.add_picture(img_case_bytes, pic_left, pic_top, width=pic_width_inch, height=pic_height_inch)


    # --- NEW Slide 8: Parts Deep Dive (Quantity) ---
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Parts Deep Dive: Top 10 by Quantity"
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)

    # Add Parts Quantity chart (large)
    img_parts_qty_bytes = io.BytesIO(fig_parts_qty.to_image(format="png", width=1200, height=900, scale=3))

    pic_width_inch = Inches(20) # A bit narrower
    pic_height_inch = pic_width_inch * (900 / 1200)
    pic_left = (prs.slide_width - pic_width_inch) / 2
    pic_top = (prs.slide_height - pic_height_inch) / 2 + Inches(0.5)

    slide.shapes.add_picture(img_parts_qty_bytes, pic_left, pic_top, width=pic_width_inch, height=pic_height_inch)


    # --- Slide 9: Parts Deep Dive (Cost) (MODIFIED for Talking Points) ---
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Parts Deep Dive: Top 10 by Cost"
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)

    # Add Parts Cost chart (left)
    img_parts_cost_bytes = io.BytesIO(fig_parts_cost.to_image(format="png", width=1200, height=900, scale=3))
    
    pic_width_inch = Inches(18.0) # Chart on the left
    pic_height_inch = pic_width_inch * (900 / 1200)
    pic_left = Inches(0.5)
    pic_top = Inches(2.5) # Lowered to fit title
    
    slide.shapes.add_picture(img_parts_cost_bytes, pic_left, pic_top, width=pic_width_inch, height=pic_height_inch)

    # --- NEW: Talking Points for Parts Slide ---
    bullets_box_left = Inches(0.5) + pic_width_inch + Inches(0.5)
    bullets_box_width = Inches(prs.slide_width - bullets_box_left - Inches(0.5))
    
    textbox = slide.shapes.add_textbox(bullets_box_left, pic_top, bullets_box_width, Inches(7.0))
    tf = textbox.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "High-Risk Component Absorption"
    p.font.name = font_name
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(54, 164, 179)

    if summary_top_parts:
        first_part = summary_top_parts[0]
        tf.add_paragraph().text = f"Top Risk Component: The **'{first_part['item']}'** component was replaced **{first_part['qty']:.0f} time(s)**, representing a gross cost of **${first_part['gross_cost']:,.2f}** absorbed by your contract."
        tf.paragraphs[-1].font.size = Pt(18)
        tf.paragraphs[-1].level = 1

        tf.add_paragraph().text = f"Unscheduled Risk: The Top 5 parts replaced collectively represent **${sum(p['gross_cost'] for p in summary_top_parts):,.2f}** in potential unbudgeted capital expense."
        tf.paragraphs[-1].font.size = Pt(18)
        tf.paragraphs[-1].level = 1
    else:
        tf.add_paragraph().text = "No high-cost parts were used in this period, indicating excellent component health."
        tf.paragraphs[-1].font.size = Pt(18)
        tf.paragraphs[-1].level = 1

    tf.add_paragraph().text = "Proactive Planning: This data confirms the service agreement is successfully converting critical component failures into zero-cost operational events."
    tf.paragraphs[-1].font.size = Pt(18)
    tf.paragraphs[-1].font.name = font_name
    tf.paragraphs[-1].level = 1
    # --- END NEW: Talking Points for Parts Slide ---


    # --- Slide 10: Activity Analysis (MODIFIED for Talking Points) ---
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Activity Analysis: Time Spent & Recurring Issues"
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)

    img_activity_bytes = io.BytesIO(fig_activity.to_image(format="png", width=1200, height=900, scale=3))
    # Place pie chart on the left
    pic_width_inch = Inches(14)
    pic_height_inch = pic_width_inch * (900/1200)
    pic_left = Inches(0.5)
    pic_top = (prs.slide_height - pic_height_inch) / 2 + Inches(0.5)
    slide.shapes.add_picture(img_activity_bytes, pic_left, pic_top, width=pic_width_inch, height=pic_height_inch)

    # --- NEW: Talking Points for Activity Slide (Corrective Action) ---
    bullets_box_left = Inches(0.5) + pic_width_inch + Inches(0.5)
    bullets_box_width = Inches(prs.slide_width - bullets_box_left - Inches(0.5))
    
    textbox = slide.shapes.add_textbox(bullets_box_left, pic_top + Inches(1.0), bullets_box_width, Inches(7.0)) # Adjusted top
    tf = textbox.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Proactive Focus: Top Recurring Issues"
    p.font.name = font_name
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 96, 96)

    if summary_top_issues:
        first_issue = summary_top_issues[0]
        tf.add_paragraph().text = f"Top Repeat Issue: **'{first_issue['corrective_action']}'** occurred **{first_issue['Total_Occurrences']:.0f} times**, indicating a trend we can address."
        tf.paragraphs[-1].font.size = Pt(18)
        tf.paragraphs[-1].level = 1
        
        tf.add_paragraph().text = f"Target for Improvement: This issue consumed **${first_issue['Total_Cost_Absorbed']:,.2f}** of covered cost and resource time this quarter."
        tf.paragraphs[-1].font.size = Pt(18)
        tf.paragraphs[-1].level = 1
    else:
        tf.add_paragraph().text = "No recurring issues detected, indicating strong system and operational health."
        tf.paragraphs[-1].font.size = Pt(18)
        tf.paragraphs[-1].level = 1
    
    tf.add_paragraph().text = "Next Steps: Let’s collaborate on **operator training** or **procedure review** to eliminate this recurring issue next quarter, maximizing machine availability."
    tf.paragraphs[-1].font.size = Pt(18)
    tf.paragraphs[-1].font.name = font_name
    tf.paragraphs[-1].level = 1
    # --- END NEW: Talking Points for Activity Slide ---


    # --- Save presentation to memory ---
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io.getvalue()
# ----------------------------------------------

# --- NEW: AI Narrative Generation Function ---
def get_ai_narrative(kpi_data: dict, model_name: str):
    """
    Generates a QBR narrative focused on validating the value and benefits 
    the customer received from their existing service agreement.
    """
    try:
        # Convert data to a JSON string for the prompt
        data_str = json.dumps(kpi_data, indent=2)

        # This prompt focuses on reporting contract value and proactive planning,
        # with strict instructions to use NO special formatting.
        prompt = f"""
You are an expert service analyst leading a Quarterly Service Review (QBR) for a valued customer who ALREADY has a full service agreement.
Your goal is to use the provided data to demonstrate the tangible value, risk mitigation, and operational continuity the agreement delivered during the last quarter/period.

Your audience consists of:
1.  Administrators (Budget/Risk Owners)
2.  Oncologists/Physicians (Throughput Owners)
3.  Physicists/Service Users (Technical Owners)

Here is the summary of the service activity and financial metrics collected under the service agreement:
{data_str}

**Instructions (Strict Output Format):**
1.  **DO NOT USE ANY MARKDOWN SYNTAX.** This includes asterisks (*), bolding (**), hyphens for lists (-), or number signs (#). Use line breaks only.
2.  All sections must be clear headings followed by brief, concise sentences.
3.  **Executive Summary:** Start by confirming the total value delivered—specifically, contrast the high Gross Cost the customer avoided with the fixed cost they paid. Conclude that the agreement performed exactly as designed.
4.  **For the Administrator (Financial Shield):**
    * Financial Protection: Use the Total Cost and high-cost anomalies to illustrate the large, volatile expenses that were successfully absorbed by the fixed contract fee. State clearly: "You exchanged unpredictable risk for budget certainty."
    * Efficiency: Use the Total Service Events and Total Labor Hours to show how frequently the contract provided "free" support, translating high activity into high value delivered.
5.  **For the Oncologist/Physician (Uptime Guarantee):**
    * Proactive Care: Focus on the Total Service Events and the Total Labor Hours. Explain that the high number of Service Visits and hours represents successful interventions that minimized machine downtime, directly ensuring maximum patient throughput.
    * Response Efficiency: Focus on the Avg Labor Cost per Service Visit as a proxy for the quick, efficient response time facilitated by the contract (no need for purchase order delays).
6.  **For the Physicist/Service User (Technical Risk Mitigation):**
    * Risk Management: Highlight the Top 5 Expensive Parts replaced. Explain that replacing these high-value, critical components under the contract successfully mitigated major technical failure risks at zero variable cost.
    * Next Steps: Use the recurring issue data to identify one area for proactive planning (e.g., operator training or PM scheduling optimization) to reduce future reactive events even further.
7.  **Conclusion:** Summarize the value and transition the conversation to planning for the next period.
8.  Ensure the entire output is structured using only plain text and new lines.
"""

        # Use streaming for a better user experience
        response_stream = ollama.chat(
            model=model_name,
            messages=[{'role': 'user', 'content': prompt}],
            stream=True
        )

        return response_stream

    except ollama.ResponseError as e:
        error_message = f"""
        ❌ Ollama Error:
        The model {model_name} likely does not exist or Ollama failed to process the request.
        1. Check model name: Ensure {model_name} is a model you have pulled (e.g., ollama pull {model_name}).
        2. Check the model's capabilities: The model may not be capable of processing the detailed instructions. Try llama3.
        
        Details: {e}
        """
        def error_generator():
            yield error_message
        return error_generator()
        
    except Exception as e:
        # Handle connection errors if Ollama isn't running
        error_message = f"""
        ❌ Connection Error: Could not connect to Ollama.

        Please ensure Ollama is running on your local machine by opening a terminal and running:
        ollama serve

        Details: {e}
        """
        # Return a generator that yields the error message
        def error_generator():
            yield error_message
        return error_generator()
# -------------------------------------------
# --- CONFIGURATION AND SIDEBAR SETUP ---

# Use a specific page title for clarity in a multi-page app
st.set_page_config(layout="wide", page_title="Labor & Service Analysis")

st.sidebar.title('🛠️ Settings')
st.sidebar.divider()

# Placeholder URL for a Labor Report - adjust this to your actual system
LABOR_URL = 'https://elekta.my.salesforce.com/00OKf000000ZH75'

# Sidebar file uploader
uploaded_file = st.sidebar.file_uploader(label='Load Service Report File (Labor & Parts):', key='labor', type=['xlsx', 'xls', 'csv'])

# --- INSTRUCTION SECTION (if no file is uploaded) ---

if uploaded_file is None:
    # --- UI/UX IMPROVEMENT: Cleaned up landing page ---
    st.markdown("<h1 style='text-align: center; color: rgb(43, 101, 124);'>📈 Labor & Service Analysis</h1>", unsafe_allow_html=True)
    st.markdown("<p style='text-align: center;'>Please upload your service report to begin.</p>", unsafe_allow_html=True)
    
    st.info(
        f"""
        **Instructions**
        1.  Go to Reports on CLM and select your **Labor/Service Activity Report**.
            * *CLM Labor Report: ({LABOR_URL})*
        2.  Select **Export** > **Details Only** > Format .csv or .xlsx.
        3.  Upload your file using the sidebar on the left.
        4.  Review the generated **Labor & Cost Analysis**.
        """
    )
    # --- END UI/UX IMPROVEMENT ---

# --- DATA PROCESSING AND ANALYSIS ---

if uploaded_file is not None:
    # File reading logic
    df = None
    file_extension = uploaded_file.name.split('.')[-1]

    if file_extension in ['xlsx', 'xls']:
        try:
            df = pd.read_excel(uploaded_file)
        except Exception as e:
            st.error(f"Error reading Excel file: {e}")
            df = None
    elif file_extension == 'csv':
        try:
            df = pd.read_csv(uploaded_file, encoding='utf-8')
        except UnicodeDecodeError:
            df = pd.read_csv(uploaded_file, encoding='ISO-8859-1')
        except Exception as e:
            st.error(f"Error reading CSV file: {e}")
            df = None

    if df is not None:

        try:
            # --- Data Cleaning and Preparation ---

            # Rename columns
            df.rename(columns={
                'Work Order: Work Order Number': 'work_order',
                'Work Order: Created Date': 'created_date',
                'Case Number': 'case_number', # <-- ADDED CASE NUMBER
                'Line Qty': 'qty',
                'Total Line Price': 'total_cost',
                'Line Price Per Unit': 'line_price_per_unit',
                'Discount %': 'discount_percent',
                'Technician': 'technician',
                'Location': 'location',
                'Activity Type': 'activity_type',
                'Corrective Action': 'corrective_action',
                'Order Type': 'order_type',
                'Line Type': 'line_type',
                'Item': 'item',
                'Work Order: Record Type': 'record_type' # <-- NEWLY ADDED
            }, inplace=True)

            full_df = df.copy()

            # Clean up the new case_number column
            full_df['case_number'] = full_df['case_number'].fillna('Unspecified').astype(str)

            # Robust date conversions
            full_df['created_date'] = pd.to_datetime(full_df['created_date'], errors='coerce')

            # Drop rows where date could not be parsed
            full_df.dropna(subset=['created_date'], inplace=True)

            # --- DYNAMIC DATE AND SELECTION FILTERS ---
            st.sidebar.divider()
            st.sidebar.header("📊 Data Filters")

            # Determine the min/max date of the uploaded data
            if not full_df.empty:
                min_data_date = full_df['created_date'].min().date()
                max_data_date = full_df['created_date'].max().date()

                # FIX: Ensure max_data_date is never in the future
                if max_data_date > date.today():
                    max_data_date = date.today()
            else:
                # Fallback if no data is present
                min_data_date = date.today() - timedelta(days=365)
                max_data_date = date.today()


            st.sidebar.subheader("📅 Timeframe Selection")
            today = date.today()

            # --- Last X Months Radio Button ---
            period_selection = st.sidebar.radio(
                "Quick Select:",
                ('Last 12 Months', 'Last 6 Months', 'Last 3 Months', 'Custom Range'),
                index=0
            )

            # 1. Determine the maximum valid end date for the default value
            valid_default_end_date = min(today, max_data_date)


            if period_selection == 'Last 12 Months':
                default_start_date = valid_default_end_date - timedelta(days=365)
                selected_end_date = valid_default_end_date
            elif period_selection == 'Last 6 Months':
                default_start_date = valid_default_end_date - timedelta(days=182)
                selected_end_date = valid_default_end_date
            elif period_selection == 'Last 3 Months':
                default_start_date = valid_default_end_date - timedelta(days=91)
                selected_end_date = valid_default_end_date
            elif period_selection == 'Custom Range':
                # For custom, default the start to the earliest possible data date
                default_start_date = min_data_date
                selected_end_date = max_data_date

            # 2. Adjust default start date: it can't be earlier than min_data_date
            default_start_date = max(default_start_date, min_data_date)


            # --- Date Picker Inputs ---

            start_date_input = st.sidebar.date_input(
                "Start Date:",
                value=default_start_date,
                min_value=min_data_date,
                max_value=max_data_date,
                key='start_date_input'
            )

            end_date_input = st.sidebar.date_input(
                "End Date:",
                value=selected_end_date,
                min_value=start_date_input,
                max_value=max_data_date,
                key='end_date_input'
            )

            # Convert date_input to datetime for comparison
            START_DATE_FILTER = pd.to_datetime(start_date_input)
            END_DATE_FILTER = pd.to_datetime(end_date_input) + timedelta(days=1)

            st.sidebar.divider()

            # --- Location Filter ---
            st.sidebar.subheader("Location")
            full_df['location'] = full_df['location'].fillna('Unspecified').astype(str)
            all_locations = sorted(full_df['location'].unique())
            selected_locations = st.sidebar.multiselect(
                "Select Locations",
                options=all_locations,
                default=all_locations
            )

            # --- Technician Filter ---
            st.sidebar.subheader("Technician")
            full_df['technician'] = full_df['technician'].fillna('Unspecified').astype(str)
            all_techs = sorted(full_df['technician'].unique())
            selected_techs = st.sidebar.multiselect(
                "Select Technicians",
                options=all_techs,
                default=all_techs
            )
            
            # --- NEW: Record Type Filter ---
            st.sidebar.subheader("Work Order Record Type")
            # Check if the column was successfully renamed/exists
            if 'record_type' in full_df.columns:
                full_df['record_type'] = full_df['record_type'].fillna('Unspecified').astype(str)
                all_record_types = sorted(full_df['record_type'].unique())
                selected_record_types = st.sidebar.multiselect(
                    "Select Record Types",
                    options=all_record_types,
                    default=all_record_types
                )
            else:
                st.sidebar.warning("Column 'Work Order: Record Type' not found in file.")
                # Create a placeholder to avoid errors
                selected_record_types = []
                all_record_types = []
            # --- END NEW FILTER ---
            
            st.sidebar.divider()
            
            # --- UI/UX IMPROVEMENT: Moved Global Toggles to Sidebar ---
            st.sidebar.subheader("📊 Global Dashboard Settings")
            include_parts = st.sidebar.toggle("Include Parts Cost in Total Cost Metrics", value=True)
            include_discounts = st.sidebar.toggle("Apply Discounts (Show Net Cost)", value=True)
            st.sidebar.caption(f"**Cost Basis:** {'Net Cost (Actual Paid)' if include_discounts else 'Gross Cost (Full Price)'}")
            st.sidebar.divider()
            # --- END UI/UX IMPROVEMENT ---
            
            # --- NUMERICAL CONVERSION AND DEBUGGING LOGIC ---
            # --- UI/UX IMPROVEMENT: Collapsed Data Integrity Check ---
            with st.sidebar.expander("✅ View Data Integrity Check", expanded=False):
                # 1. Total Cost Check
                full_df['total_cost_temp'] = pd.to_numeric(full_df['total_cost'], errors='coerce')
                non_numeric_cost = full_df[full_df['total_cost_temp'].isna()].copy()

                # 2. Quantity Check
                full_df['qty_temp'] = pd.to_numeric(full_df['qty'], errors='coerce')
                non_numeric_qty = full_df[full_df['qty_temp'].isna()].copy()

                # 3. Line Price Per Unit Check
                full_df['line_price_per_unit_temp'] = pd.to_numeric(full_df['line_price_per_unit'], errors='coerce')
                non_numeric_ppu = full_df[full_df['line_price_per_unit_temp'].isna()].copy()

                # 4. Discount % Check
                full_df['discount_percent_temp'] = pd.to_numeric(full_df['discount_percent'], errors='coerce')
                non_numeric_discount = full_df[full_df['discount_percent_temp'].isna()].copy()


                # Final assignment (all non-numeric and NaN values become 0)
                full_df['total_cost'] = full_df['total_cost_temp'].fillna(0)
                full_df['qty'] = full_df['qty_temp'].fillna(0)
                full_df['line_price_per_unit'] = full_df['line_price_per_unit_temp'].fillna(0)
                full_df['discount_percent'] = full_df['discount_percent_temp'].fillna(0)

                # Display non-numeric results
                if not non_numeric_cost.empty:
                    st.warning(f"⚠️ **{len(non_numeric_cost)}** Non-Numeric Total Cost Entries Found:")
                    st.dataframe(non_numeric_cost[['work_order', 'line_type', 'total_cost', 'qty']], use_container_width=True)

                if not non_numeric_qty.empty:
                    st.warning(f"⚠️ **{len(non_numeric_qty)}** Non-Numeric Qty/Hour Entries Found:")
                    st.dataframe(non_numeric_qty[['work_order', 'line_type', 'total_cost', 'qty']], use_container_width=True)

                if not non_numeric_ppu.empty:
                    st.warning(f"⚠️ **{len(non_numeric_ppu)}** Non-Numeric Price Per Unit Entries Found:")
                    st.dataframe(non_numeric_ppu[['work_order', 'line_type', 'line_price_per_unit']], use_container_width=True)

                if not non_numeric_discount.empty:
                    st.warning(f"⚠️ **{len(non_numeric_discount)}** Non-Numeric Discount % Entries Found:")
                    st.dataframe(non_numeric_discount[['work_order', 'line_type', 'discount_percent']], use_container_width=True)
            # --- END UI/UX IMPROVEMENT ---


            # Clean up and fill missing order types
            full_df['order_type'] = full_df['order_type'].fillna('Unspecified').astype(str)
            full_df.dropna(subset=['work_order'], inplace=True)


            # --- APPLY ALL FILTERS TO THE MAIN DATAFRAME ---
            # Create the record type filter mask
            if 'record_type' in full_df.columns:
                record_type_mask = full_df['record_type'].isin(selected_record_types)
            else:
                # If column doesn't exist, create a mask that includes everything
                record_type_mask = pd.Series(True, index=full_df.index)

            full_df = full_df[
                (full_df['created_date'] >= START_DATE_FILTER) &
                (full_df['created_date'] < END_DATE_FILTER) &
                (full_df['location'].isin(selected_locations)) &
                (full_df['technician'].isin(selected_techs)) &
                (record_type_mask) # <-- NEW FILTER APPLIED
            ].copy()

            # Check if any data remains after filtering
            if full_df.empty:
                 st.error("No data remains after applying the filters. Please adjust your date range, locations, or technicians.")
                 st.stop() # FIX: Used st.stop() to halt Streamlit script execution

            # Define Line Type for filtering
            if 'line_type' in full_df.columns:
                 full_df['line_type'] = full_df['line_type'].astype(str).str.lower()

                 # 1. Identify Labor Lines (Labor total_cost is used as-is for labor, qty is hours)
                 labor_filter = full_df['line_type'].str.contains('labor|time|service', na=False)
                 labor_df = full_df[labor_filter].copy()
                 labor_df['labor_hours'] = labor_df['qty']

                 # Labor Gross Cost: We assume the original 'total_cost' (Total Line Price) is the Gross price.
                 labor_df['labor_gross_cost'] = labor_df['total_cost']

                 # 2. Identify Parts Lines
                 parts_filter = ~labor_filter
                 parts_df = full_df[parts_filter].copy()
                 parts_df.dropna(subset=['item'], inplace=True)

                 # Parts Gross Cost: Calculated from PPU * QTY
                 parts_df['parts_gross_cost'] = parts_df['line_price_per_unit'] * parts_df['qty']

            else:
                 st.warning("Warning: 'Line Type' column not found. Cannot separate labor and parts costs.")
                 labor_df = full_df.copy()
                 parts_df = pd.DataFrame()
                 labor_df['labor_hours'] = labor_df['qty']


            # Drop rows where essential data is missing for time-based analysis
            labor_df.dropna(subset=['work_order'], inplace=True)
            parts_df.dropna(subset=['work_order'], inplace=True)


            # --- VISUALIZATION SECTION ---

            st.markdown("<h1 style='color: rgb(43, 101, 124);'>📈 Labor & Parts Service Analysis</h1>", unsafe_allow_html=True)
            # Note: This info box is not updated to include record types, but it could be
            st.info(f"Displaying data from **{START_DATE_FILTER.strftime('%Y-%m-%d')}** to **{(END_DATE_FILTER - timedelta(days=1)).strftime('%Y-%m-%d')}** for **{len(selected_locations)}** location(s), **{len(selected_techs)}** technician(s), and **{len(selected_record_types)}** record type(s).")

            # --- START: GLOBAL TOGGLES AND COST CALCULATION ---
            
            # --- CONDITIONAL COST ASSIGNMENT BASED ON DISCOUNT TOGGLE ---

            # Calculate line-specific discount factors (1 - Discount %)
            labor_discount_factor_line = 1 - (labor_df['discount_percent'] / 100)
            parts_discount_factor_line = 1 - (parts_df['discount_percent'] / 100)

            # Calculate true Net Cost columns
            labor_df['labor_net_cost_calc'] = labor_df['labor_gross_cost'] * labor_discount_factor_line
            parts_df['parts_net_cost_calc'] = parts_df['parts_gross_cost'] * parts_discount_factor_line

            # Use the toggle to decide which cost basis to use for 'total_cost'
            if include_discounts:
                # Assign Net Cost
                labor_df['total_cost'] = labor_df['labor_net_cost_calc']
                parts_df['total_cost'] = parts_df['parts_net_cost_calc']
                parts_label_suffix = "Net"
            else:
                # Assign Gross Cost
                labor_df['total_cost'] = labor_df['labor_gross_cost']
                parts_df['total_cost'] = parts_df['parts_gross_cost']
                parts_label_suffix = "Gross"

            # Ensure no negative costs result from calculation
            labor_df['total_cost'] = np.maximum(labor_df['total_cost'], 0)
            parts_df['total_cost'] = np.maximum(parts_df['total_cost'], 0)

            # --- CREATE THE FINAL, TOGGLE-AWARE DATAFRAME ---
            # This df is now available to ALL tabs
            full_df_filtered = pd.concat([labor_df, parts_df], ignore_index=True)

            # --- END: GLOBAL TOGGLES AND COST CALCULATION ---

            
            # --- START: PRE-CALCULATE ALL KPIs FOR PPT AND AI ---
            # This section calculates all KPIs needed for both the Streamlit UI and the PowerPoint/AI
            
            # --- Main KPI Calculations (from tab_kpi) ---
            total_labor_current_cost = labor_df['total_cost'].sum()
            total_parts_current_cost = parts_df['total_cost'].sum()
            total_labor_gross_cost = labor_df['labor_gross_cost'].sum()
            total_parts_gross_cost = parts_df['parts_gross_cost'].sum()
            total_parts_replaced = parts_df['qty'].sum()
            total_hours = labor_df['labor_hours'].sum()

            date_grouper = full_df_filtered['created_date'].dt.date.rename('event_date')
            df_wo_count = full_df_filtered.groupby(['work_order', date_grouper]).agg(
                wo_count=pd.NamedAgg(column='work_order', aggfunc='first'),
                created_date=pd.NamedAgg(column='created_date', aggfunc='first')
            ).reset_index()
            # Renaming the unit of measure for consistency in KPIs
            total_events = df_wo_count['work_order'].nunique() 
            event_noun = "Service Visit (WO)" # <-- NEW: Consistent Terminology

            if include_parts:
                total_tcs = total_labor_current_cost + total_parts_current_cost
                tcs_label = f"Total Cost (Labor + Parts {parts_label_suffix})"
            else:
                total_tcs = total_labor_current_cost
                tcs_label = f"Total Cost (Labor {parts_label_suffix} Only)"

            # --- NEW: Calculate Average Costs ---
            # 1. Avg Total Cost per Event (Service Visit)
            avg_tcs_per_event = total_tcs / total_events if total_events > 0 else 0
            
            # 2. Avg Labor Only Cost per Event (Service Visit)
            avg_labor_only_per_event = total_labor_current_cost / total_events if total_events > 0 else 0


            if include_parts:
                base_df_for_stats = full_df_filtered
            else:
                base_df_for_stats = labor_df
            df_event_costs = base_df_for_stats.groupby('work_order')['total_cost'].sum()

            if not df_event_costs.empty:
                median_cost = df_event_costs.median()
                max_cost = df_event_costs.max()
                min_cost = df_event_costs.min()
            else:
                median_cost = 0
                max_cost = 0
                min_cost = 0

            total_labor_net_cost = labor_df['labor_net_cost_calc'].sum()
            total_parts_net_cost = parts_df['parts_net_cost_calc'].sum()
            total_discount_given = (total_labor_gross_cost - total_labor_net_cost) + (total_parts_gross_cost - total_parts_net_cost)

            # --- Case KPI Calculations (from tab_case) ---
            NON_CASE_ORDER_TYPES = ['Preventive Maintenance', 'FCO', 'Unspecified']
            case_df = full_df_filtered[
                ~full_df_filtered['order_type'].isin(NON_CASE_ORDER_TYPES) &
                (full_df_filtered['case_number'] != 'Unspecified')
            ].copy()

            total_cases = 0
            avg_cost_per_case = 0
            avg_visits_per_case = 0
            median_case_cost = 0
            max_case_cost = 0
            min_case_cost = 0
            df_case_agg = pd.DataFrame() # Initialize empty

            if not case_df.empty:
                case_df['parts_cost'] = np.where(case_df['labor_hours'].isna(), case_df['total_cost'], 0)
                case_df['labor_cost'] = np.where(case_df['labor_hours'].notna(), case_df['total_cost'], 0)
                
                df_case_agg = case_df.groupby('case_number').agg(
                    total_cost_per_case=pd.NamedAgg(column='total_cost', aggfunc='sum'),
                    parts_cost_per_case=pd.NamedAgg(column='parts_cost', aggfunc='sum'),
                    labor_cost_per_case=pd.NamedAgg(column='labor_cost', aggfunc='sum'),
                    total_hours_per_case=pd.NamedAgg(column='labor_hours', aggfunc='sum'),
                    visits_per_case=pd.NamedAgg(column='work_order', aggfunc='nunique'),
                    first_visit_date=pd.NamedAgg(column='created_date', aggfunc='min')
                ).reset_index().sort_values(by='total_cost_per_case', ascending=False)

                if not df_case_agg.empty:
                    total_cases = df_case_agg['case_number'].nunique()
                    avg_cost_per_case = df_case_agg['total_cost_per_case'].mean()
                    avg_visits_per_case = df_case_agg['visits_per_case'].mean()
                    median_case_cost = df_case_agg['total_cost_per_case'].median()
                    max_case_cost = df_case_agg['total_cost_per_case'].max()
                    min_case_cost = df_case_agg['total_cost_per_case'].min()

            # --- Corrective Action Summary for PPT (Must be calculated outside of tab_activity for PPT button) ---
            df_corrective_analysis = full_df_filtered.groupby(['location', 'corrective_action']).agg(
                occurrence_count=pd.NamedAgg(column='work_order', aggfunc='nunique'),
                total_cost_for_action=pd.NamedAgg(column='total_cost', aggfunc='sum')
            ).reset_index()

            # Sum the occurrences across all locations to find the top issues
            df_issue_summary = df_corrective_analysis.groupby('corrective_action').agg(
                Total_Occurrences=pd.NamedAgg(column='occurrence_count', aggfunc='sum'),
                Total_Cost_Absorbed=pd.NamedAgg(column='total_cost_for_action', aggfunc='sum')
            ).reset_index().sort_values(by='Total_Occurrences', ascending=False)
            
            # Filter out 'Unspecified' and only show the top 10
            df_issue_summary = df_issue_summary[
                ~df_issue_summary['corrective_action'].astype(str).str.contains('unspecified|n/a', case=False, na=False)
            ].head(10)
            
            # --- END: PRE-CALCULATE ALL KPIs ---


            # --- START: TABS ---
            tab_kpi, tab_performance, tab_activity, tab_parts, tab_case, tab_data, tab_ai_narrative = st.tabs([ # <-- ADDED tab_ai_narrative
                "🎯 KPI Dashboard",
                "👨‍🔧 Performance",
                "⚡ Activity",
                "📦 Parts Deep Dive",
                "🕵️ Case Analysis",
                "Raw Data",
                "💡 AI Narrative" # <-- NEW TAB NAME
            ])

            # --- Initialize Figure Variables ---
            # These need to be defined *before* the try block ends so the PPT button can access them
            fig_trend = go.Figure()
            fig_split = go.Figure()
            fig_tech = go.Figure()
            fig_loc = go.Figure()
            fig_activity = go.Figure()
            fig_parts_qty = go.Figure()
            fig_parts_cost = go.Figure()
            # --- NEW: Define both case figures ---
            fig_case_trend_total = go.Figure()
            fig_case_heatmap = go.Figure()

            with tab_kpi:
                # 1. Overall Metrics and Trend
                with st.container(border=True):

                    st.subheader("🎯 Service Insights")
                    st.caption("Metrics below reflect the global settings selected in the sidebar.")
                    st.divider()

                    # --- RE-AGGREGATE MONTHLY DATA AFTER TOGGLE DECISION ---

                    # Group labor costs and hours monthly (Labor_Cost is now conditional: Net or Gross)
                    df_labor_monthly = labor_df.groupby(labor_df['created_date'].dt.to_period('M')).agg(
                        Labor_Cost=pd.NamedAgg(column='total_cost', aggfunc='sum'), # Current Cost (Toggled)
                        Labor_Cost_Gross=pd.NamedAgg(column='labor_gross_cost', aggfunc='sum'),
                        Total_Hours=pd.NamedAgg(column='labor_hours', aggfunc='sum')
                    ).reset_index()

                    # Group parts cost monthly (Parts_Cost is now conditional: Net or Gross)
                    df_parts_monthly = parts_df.groupby(parts_df['created_date'].dt.to_period('M')).agg(
                        Parts_Cost=pd.NamedAgg(column='total_cost', aggfunc='sum'), # Current Cost (Toggled)
                        Parts_Cost_Gross=pd.NamedAgg(column='parts_gross_cost', aggfunc='sum')
                    ).reset_index()

                    # Group parts quantity (qty) monthly
                    df_parts_qty_monthly = parts_df.groupby(parts_df['created_date'].dt.to_period('M')).agg(
                        Total_Parts_Qty=pd.NamedAgg(column='qty', aggfunc='sum')
                    ).reset_index()

                    # Group WO count monthly (uses df_wo_count from KPI pre-calc)
                    df_wo_monthly = df_wo_count.groupby(df_wo_count['created_date'].dt.to_period('M')).agg(
                        Total_Events=pd.NamedAgg(column='work_order', aggfunc='nunique')
                    ).reset_index()


                    # --- MERGE ALL MONTHLY DATA ---
                    df_monthly_combined = pd.merge(df_labor_monthly, df_parts_monthly, on='created_date', how='outer').fillna(0)
                    df_monthly_combined = pd.merge(df_monthly_combined, df_parts_qty_monthly, on='created_date', how='outer').fillna(0)
                    df_monthly_combined = pd.merge(df_monthly_combined, df_wo_monthly, on='created_date', how='outer').fillna(0)
                    df_monthly_combined['Total_Cost'] = df_monthly_combined['Labor_Cost'] + df_monthly_combined['Parts_Cost']
                    
                    # --- FIX: Create string X-axis for plotting ---
                    df_monthly_combined['created_date'] = df_monthly_combined['created_date'].dt.to_timestamp()
                    # Sort by the actual timestamp first
                    df_monthly_combined = df_monthly_combined.sort_values('created_date')
                    # Create a string column for plotting
                    df_monthly_combined['Month_Str'] = df_monthly_combined['created_date'].dt.strftime('%b-%Y')
                    # --- END FIX ---

                    # --- KPI LAYOUT (Modified) ---
                    # Uses pre-calculated KPIs from the global scope

                    # Row 1: Core Costs and Visits
                    col1, col2, col3, col4 = st.columns(4)
                    with col1: st.metric(label=tcs_label, value=f"${total_tcs:,.2f}")
                    with col2:
                        st.metric(label=f"Labor Cost ({parts_label_suffix})", value=f"${total_labor_current_cost:,.2f}")
                        context_cost_labor = total_labor_gross_cost if include_discounts else total_labor_net_cost
                        context_label_labor = "Gross" if include_discounts else "Net"
                        st.markdown(f"<p style='font-size: 12px; color: gray;'>{context_label_labor}: ${context_cost_labor:,.2f}</p>", unsafe_allow_html=True)

                    with col3:
                        st.metric(label=f"Parts Cost ({parts_label_suffix})", value=f"${total_parts_current_cost:,.2f}")
                        context_cost_parts = total_parts_gross_cost if include_discounts else total_parts_net_cost
                        context_label_parts = "Gross" if include_discounts else "Net"
                        st.markdown(f"<p style='font-size: 12px; color: gray;'>{context_label_parts}: ${context_cost_parts:,.2f}</p>", unsafe_allow_html=True)

                    with col4: st.metric(label="Total Service Visits (WOs)", value=f"{total_events:,}") # <-- RENAMED

                    # Row 2: Efficiency, Volume, and Discount Metrics
                    if include_discounts:
                        col5, col6, col7, col8 = st.columns(4)
                        with col5:
                            st.metric(label=f"Avg Total Cost per {event_noun}", value=f"${avg_tcs_per_event:,.2f}") # <-- RENAMED
                            st.markdown(f"<p style='font-size: 12px; color: gray;'>Med: ${median_cost:,.2f} | Max: ${max_cost:,.2f} | Min: ${min_cost:,.2f}</p>", unsafe_allow_html=True)
                        with col6: 
                            st.metric(label=f"Avg Labor Cost per {event_noun}", value=f"${avg_labor_only_per_event:,.2f}") # <-- NEW METRIC
                            st.markdown(f"<p style='font-size: 12px; color: gray;'>Total Hours: {total_hours:,.1f} h</p>", unsafe_allow_html=True)
                        with col7: st.metric(label="Total Parts Replaced (Qty)", value=f"{total_parts_replaced:,.0f}")
                        with col8: st.metric(label="Total Discount Given", value=f"${total_discount_given:,.2f}")
                    else:
                        col5, col6, col7 = st.columns(3)
                        with col5:
                            st.metric(label=f"Avg Total Cost per {event_noun}", value=f"${avg_tcs_per_event:,.2f}") # <-- RENAMED
                            st.markdown(f"<p style='font-size: 12px; color: gray;'>Med: ${median_cost:,.2f} | Max: ${max_cost:,.2f} | Min: ${min_cost:,.2f}</p>", unsafe_allow_html=True)
                        with col6: 
                            st.metric(label=f"Avg Labor Cost per {event_noun}", value=f"${avg_labor_only_per_event:,.2f}") # <-- NEW METRIC
                            st.markdown(f"<p style='font-size: 12px; color: gray;'>Total Hours: {total_hours:,.1f} h</p>", unsafe_allow_html=True)
                        with col7: st.metric(label="Total Parts Replaced (Qty)", value=f"{total_parts_replaced:,.0f}")

                    st.divider()

                    # --- PARTS TREND CHART ---
                    st.markdown("### Monthly Parts Analysis Trend")

                    parts_metric_options = ['Parts_Cost', 'Total_Parts_Qty', 'Parts_Cost_Gross']

                    parts_metric = st.radio(
                        "Select Parts Metric:",
                        parts_metric_options,
                        format_func=lambda x: x.replace("_", " ").title().replace("Parts Cost", f"Parts Cost ({parts_label_suffix})") if x == 'Parts_Cost' else x.replace("_", " ").title().replace("Parts Cost Gross", "Parts Cost (Gross)"),
                        horizontal=True,
                        key='parts_metric_radio'
                    )
                    
                    # ... (Info box logic) ...
                    current_parts_cost_label = f"Total Parts Cost ({parts_label_suffix})"
                    if parts_metric == 'Parts_Cost':
                        st.info(f"The **{current_parts_cost_label}** over this period is **${total_parts_current_cost:,.2f}**.")
                    elif parts_metric == 'Parts_Cost_Gross':
                        st.info(f"The **Total Parts Cost (Gross)** over this period is **${total_parts_gross_cost:,.2f}**.")
                    else:
                        st.info(f"The **Total Parts Replaced** over this period is **{total_parts_replaced:,.0f} units**.")


                    # --- FIX: Use 'Month_Str' for X-axis ---
                    df_parts_plot = df_monthly_combined[['Month_Str', 'created_date', parts_metric]].sort_values('created_date')
                    
                    fig_parts_trend = px.area(
                        df_parts_plot,
                        x='Month_Str', # <-- CHANGED
                        y=parts_metric,
                        title=f'Monthly Trend for {parts_metric.replace("_", " ").title().replace("Parts Cost", f"Parts Cost ({parts_label_suffix})").replace("Parts Cost Gross", "Parts Cost (Gross)")}',
                        labels={parts_metric: parts_metric.replace("_", " ").title(), 'Month_Str': 'Month'}, # <-- CHANGED
                        color_discrete_sequence=['rgb(54, 164, 179)'],
                        template="streamlit"
                    )

                    fig_parts_trend.update_xaxes(type='category') # Treat as category
                    fig_parts_trend.update_layout(hovermode="x unified")
                    # FIX: use_container_width=True -> width='stretch'
                    st.plotly_chart(fig_parts_trend, use_container_width=True)

                    st.divider()

                    # --- MAIN TREND (COST, HOURS, EVENTS) ---
                    st.markdown("### Total Cost, Hours, & Service Visit Trend")

                    if include_parts:
                        metric_options = ['Total_Cost', 'Labor_Cost', 'Total_Hours', 'Total_Events']
                    else:
                        metric_options = ['Labor_Cost', 'Total_Hours', 'Total_Events']

                    selected_metric = st.radio(
                        "Select Metric to Visualize:",
                        metric_options,
                        format_func=lambda x: x.replace("_", " ").title(),
                        horizontal=True,
                        key='trend_metric_radio'
                    )
                    
                    # ... (Info box logic) ...
                    if selected_metric == 'Total_Cost':
                        st.info(f"The **{tcs_label}** over this period is **${total_tcs:,.2f}**.")
                    elif selected_metric == 'Labor_Cost':
                        st.info(f"Total Labor Cost ({parts_label_suffix}): **${total_labor_current_cost:,.2f}**.")
                    elif selected_metric == 'Total_Hours':
                        st.info(f"Total Labor Hours: **{total_hours:,.1f} h**.")
                    elif selected_metric == 'Total_Events':
                        st.info(f"Total Service Visits (WOs): **{total_events:,}**.")

                    
                    # --- FIX: Use 'Month_Str' for X-axis ---
                    df_long = pd.melt(
                        df_monthly_combined,
                        id_vars=['created_date', 'Month_Str'], # <-- CHANGED
                        value_vars=metric_options,
                        var_name='Metric',
                        value_name='Value'
                    )

                    df_filtered = df_long[df_long['Metric'] == selected_metric].sort_values('created_date')

                    # Assign to fig_trend for PowerPoint
                    fig_trend = px.area(
                        df_filtered,
                        x='Month_Str', # <-- CHANGED
                        y='Value',
                        title=f'Monthly Trend for {selected_metric.replace("_", " ").title()}',
                        labels={'Value': selected_metric.replace("_", " ").title(), 'Month_Str': 'Month'}, # <-- CHANGED
                        color_discrete_sequence=['rgb(43, 101, 125)'],
                        template="streamlit"
                    )

                    fig_trend.update_xaxes(type='category') # Treat as category
                    fig_trend.update_layout(hovermode="x unified")
                    # FIX: use_container_width=True -> width='stretch'
                    st.plotly_chart(fig_trend, use_container_width=True)

                    st.subheader("Cost Split: Labor vs. Parts")

                    if include_discounts:
                        st.info(f"**Total Labor Cost ({parts_label_suffix}):** ${total_labor_current_cost:,.2f} | **Total Parts Cost ({parts_label_suffix}):** ${total_parts_current_cost:,.2f} | **Total Discount Given:** ${total_discount_given:,.2f}")
                    else:
                        st.info(f"**Total Labor Cost ({parts_label_suffix}):** ${total_labor_current_cost:,.2f} | **Total Parts Cost ({parts_label_suffix}):** ${total_parts_current_cost:,.2f}")
                    
                    # --- FIX: Use 'Month_Str' for X-axis ---
                    df_split_melt = pd.melt(
                        df_monthly_combined.sort_values('created_date'), # Ensure sort
                        id_vars=['Month_Str'], # <-- CHANGED
                        value_vars=['Labor_Cost', 'Parts_Cost']
                    )

                    # Assign to fig_split for PowerPoint
                    fig_split = px.area(
                        df_split_melt,
                        x='Month_Str', # <-- CHANGED
                        y='value',
                        color='variable',
                        title='Monthly Breakdown of Total Service Cost (TCS)',
                        labels={'value': 'Cost ($)', 'variable': 'Cost Type', 'Month_Str': 'Month'}, # <-- CHANGED
                        color_discrete_map={'Labor_Cost': 'rgb(43, 101, 125)', 'Parts_Cost': 'rgb(54, 164, 179)'},
                        template="streamlit"
                    )
                    fig_split.update_xaxes(type='category') # Treat as category
                    # FIX: use_container_width=True -> width='stretch'
                    st.plotly_chart(fig_split, use_container_width=True)


            # 3. Technician and Location Performance
            with tab_performance:
                st.title("👨‍🔧 Technician & Location Performance")
                st.caption(f"Costs shown as **{parts_label_suffix} Cost** based on global settings.")

                tech_col, loc_col = st.columns(2)

                # --- Technician Analysis (Bar Chart) - Always Labor Only ---
                with tech_col:
                    # --- UI/UX IMPROVEMENT: Added container ---
                    with st.container(border=True):
                        st.subheader("Top 10 Technicians by Labor Hours")

                        df_tech = labor_df.groupby('technician').agg(
                            total_hours=pd.NamedAgg(column='labor_hours', aggfunc='sum'),
                            total_cost=pd.NamedAgg(column='total_cost', aggfunc='sum'),
                            wo_count=pd.NamedAgg(column='work_order', aggfunc='nunique')
                        ).sort_values(by='total_hours', ascending=True).tail(10).reset_index()

                        # Assign to fig_tech for PowerPoint
                        fig_tech = px.bar(
                            df_tech,
                            y='technician',
                            x='total_hours',
                            hover_data={'total_cost': True, 'wo_count': True},
                            title='Total Hours Booked (Labor Only)',
                            orientation='h',
                            color='total_cost',
                            color_continuous_scale=px.colors.sequential.Teal
                        )
                        # FIX: use_container_width=True -> width='stretch'
                        st.plotly_chart(fig_tech, use_container_width=True)

                # --- Location Analysis (Treemap) ---
                with loc_col:
                    # --- UI/UX IMPROVEMENT: Added container ---
                    with st.container(border=True):
                        cost_title_loc = f"Total Cost (Labor + Parts {parts_label_suffix})" if include_parts else f"Total Cost (Labor {parts_label_suffix} Only)"
                        st.subheader(f"Cost Distribution by Location ({cost_title_loc})")

                        if include_parts:
                            # Use the globally filtered, toggle-aware dataframe
                            df_loc_base = full_df_filtered
                        else:
                            df_loc_base = labor_df

                        df_loc = df_loc_base.groupby('location').agg(
                            total_cost=pd.NamedAgg(column='total_cost', aggfunc='sum'),
                        ).reset_index().sort_values(by='total_cost', ascending=False).head(10)

                        df_loc_hours = labor_df.groupby('location').agg(total_hours=pd.NamedAgg(column='labor_hours', aggfunc='sum')).reset_index()
                        df_loc = pd.merge(df_loc, df_loc_hours, on='location', how='left').fillna(0)

                        # Assign to fig_loc for PowerPoint
                        fig_loc = px.treemap(
                            df_loc,
                            path=[px.Constant("All Locations"), 'location'],
                            values='total_cost',
                            color='total_hours',
                            color_continuous_scale='Mint',
                            title=f'Top 10 Locations by {cost_title_loc}'
                        )
                        # FIX: use_container_width=True -> width='stretch'
                        st.plotly_chart(fig_loc, use_container_width=True)

            # 4. Efficiency and Activity Deep Dive
            with tab_activity:
                st.title("🔍 Efficiency and Activity Deep Dive (Labor Only)")
                st.caption(f"Costs shown as **{parts_label_suffix} Cost** based on global settings. Corrective Action analysis below includes **all** line types.")

                col_act_pie, col_issue_analysis = st.columns(2)

                # --- COLUMN 1: Activity Pie Chart (Hours) ---
                with col_act_pie:
                    # --- UI/UX IMPROVEMENT: Added container ---
                    with st.container(border=True):
                        st.subheader("Time Spent by Service Activity Type")
                        df_activity = labor_df.groupby('activity_type').agg(
                            total_hours=pd.NamedAgg(column='labor_hours', aggfunc='sum'),
                            wo_count=pd.NamedAgg(column='work_order', aggfunc='nunique')
                        ).reset_index().sort_values(by='total_hours', ascending=False).head(10)

                        # Assign to fig_activity for PowerPoint
                        fig_activity = px.pie(
                            df_activity,
                            names='activity_type',
                            values='total_hours',
                            title='Top 10 Activity Types by Hours',
                            color_discrete_sequence=px.colors.sequential.Mint_r
                        )
                        fig_activity.update_traces(textposition='inside', textinfo='percent+label')
                        st.plotly_chart(fig_activity, use_container_width=True)

                # --- COLUMN 2: Corrective Action Analysis (Top Issues) ---
                with col_issue_analysis:
                    # --- UI/UX IMPROVEMENT: Added container ---
                    with st.container(border=True):
                        st.subheader("Top Corrective Actions & Repeat Issues")
                        
                        # df_issue_summary is already calculated globally, just need to display it.

                        # Calculate Total Cost per Occurrence (Cost / WO)
                        df_corrective_analysis['avg_cost_per_wo'] = df_corrective_analysis['total_cost_for_action'] / df_corrective_analysis['occurrence_count']

                        st.write("##### Top 10 Repeated Issues (by WO Count)")
                        st.dataframe(
                            df_issue_summary, 
                            column_config={
                                "Total_Occurrences": st.column_config.NumberColumn("WO Count"),
                                "Total_Cost_Absorbed": st.column_config.NumberColumn(f"Total Cost Absorbed ({parts_label_suffix})", format="$%.2f")
                            },
                            use_container_width=True,
                            hide_index=True
                        )
                        
                        st.caption("Use these repeated issues to discuss proactive training or parts stocking with the customer.")


            # 5. Parts Deep Dive
            with tab_parts:
                st.title("📦 Parts Usage Deep Dive")
                if parts_df.empty:
                    st.warning("No parts data available after filtering.")
                else:
                    col_top_qty, col_top_cost = st.columns(2)

                    # --- Top 10 Items by Quantity (Bar Chart) ---
                    with col_top_qty:
                        # --- UI/UX IMPROVEMENT: Added container ---
                        with st.container(border=True):
                            st.subheader("Top 10 Parts by Quantity")

                            df_parts_qty = parts_df.groupby('item').agg(
                                qty=pd.NamedAgg(column='qty', aggfunc='sum'),
                                gross_cost=pd.NamedAgg(column='parts_gross_cost', aggfunc='sum')
                            ).sort_values(by='qty', ascending=True).tail(10).reset_index()

                            # Assign to fig_parts_qty for PowerPoint
                            fig_parts_qty = px.bar(
                                df_parts_qty,
                                y='item',
                                x='qty',
                                hover_data={'gross_cost': ':.2f'},
                                title='Quantity of Top 10 Parts Used',
                                orientation='h',
                                color='gross_cost', # Color by gross cost
                                color_continuous_scale=px.colors.sequential.Teal # Matching color to Labor Techs
                            )
                            # FIX: use_container_width=True -> width='stretch'
                            st.plotly_chart(fig_parts_qty, use_container_width=True)

                    # --- Top 10 Items by Gross Cost (Bar Chart) ---
                    with col_top_cost:
                        # --- UI/UX IMPROVEMENT: Added container ---
                        with st.container(border=True):
                            st.subheader(f"Top 10 Parts by Cost")

                            df_parts_cost = parts_df.groupby('item').agg(
                                gross_cost=pd.NamedAgg(column='parts_gross_cost', aggfunc='sum'),
                                qty=pd.NamedAgg(column='qty', aggfunc='sum')
                            ).sort_values(by='gross_cost', ascending=True).tail(10).reset_index()

                            # Assign to fig_parts_cost for PowerPoint
                            fig_parts_cost = px.bar(
                                df_parts_cost,
                                y='item',
                                x='gross_cost',
                                hover_data={'qty': True},
                                title='Cost of Top 10 Parts',
                                orientation='h',
                                color='qty', # Color by quantity
                                color_continuous_scale=px.colors.sequential.Mint # Matching color to Locations
                            )
                            # FIX: use_container_width=True -> width='stretch'
                            st.plotly_chart(fig_parts_cost, use_container_width=True)


            # 6. --- NEW --- Case Analysis
            with tab_case:
                st.title("🕵️ Case Level Analysis")
                st.info("This analysis excludes non-case work orders like 'Preventive Maintenance' and 'FCO' to focus on reactive service costs.")
                st.caption(f"Costs shown as **{parts_label_suffix} Cost** based on global settings.")

                # Uses 'case_df', 'df_case_agg', and KPIs from the pre-calculation step
                if case_df.empty or df_case_agg.empty:
                    st.warning("No data found for reactive service cases after filtering. This may be expected if you only filtered for PMs or your report doesn't have 'Case Number' data.")
                else:
                    # 4. --- INSIGHTS FIRST ---
                    # --- UI/UX IMPROVEMENT: Added container ---
                    with st.container(border=True):
                        st.subheader("Case-Level Insights (Reactive Service Only)")

                        # KPIs are already calculated: total_cases, avg_cost_per_case, etc.
                        
                        case_col1, case_col2, case_col3 = st.columns(3)
                        with case_col1:
                            st.metric(label="Total Cases", value=f"{total_cases:,}")
                        with case_col2:
                            cost_label_suffix = f"({parts_label_suffix})"
                            st.metric(label=f"Avg. Cost per Case {cost_label_suffix}", value=f"${avg_cost_per_case:,.2f}")
                            st.markdown(f"<p style='font-size: 12px; color: gray;'>Median: ${median_case_cost:,.2f} | Max: ${max_case_cost:,.2f} | Min: ${min_case_cost:,.2f}</p>", unsafe_allow_html=True)
                        with case_col3:
                            st.metric(label="Avg. Visits per Case", value=f"{avg_visits_per_case:,.1f}")

                    st.divider()

                    # 5. --- GRAPHS SECOND ---
                    # --- UI/UX IMPROVEMENT: Added container ---
                    with st.container(border=True):
                        st.subheader("Case Volume Trend")

                        trend_view = st.radio( 
                            "Select Trend View:",
                            ("Total Cases", "Cases by Location"), # <-- UPDATED LABEL
                            horizontal=True,
                            key='case_trend_view' # <-- This key must be unique
                        )

                        # Get the details for each case (first visit date, location)
                        df_case_details = case_df.groupby('case_number').agg(
                            first_visit_date=pd.NamedAgg(column='created_date', aggfunc='min'),
                            location=pd.NamedAgg(column='location', aggfunc='first') # Assumes one location per case
                        ).reset_index()

                        # Resample by month
                        df_case_details['visit_month'] = df_case_details['first_visit_date'].dt.to_period('M')

                        # --- Generate Trend Chart (for PPT) ---
                        df_trend_total = df_case_details.groupby('visit_month').agg(
                            case_count=pd.NamedAgg(column='case_number', aggfunc='nunique')
                        ).reset_index()
                        
                        # --- FIX: Use 'Month_Str' for X-axis ---
                        df_trend_total['visit_month_ts'] = df_trend_total['visit_month'].dt.to_timestamp()
                        df_trend_total = df_trend_total.sort_values('visit_month_ts') # Sort by timestamp
                        df_trend_total['Month_Str'] = df_trend_total['visit_month_ts'].dt.strftime('%b-%Y') # Create string

                        # Assign to fig_case_trend_total for PowerPoint
                        fig_case_trend_total = px.area(
                            df_trend_total,
                            x='Month_Str', # <-- CHANGED
                            y='case_count',
                            title='Total New Cases Over Time',
                            labels={'case_count': 'Number of Cases', 'Month_Str': 'Month'}, # <-- CHANGED
                            color_discrete_sequence=['rgb(43, 101, 125)']
                        )
                        fig_case_trend_total.update_xaxes(type='category')
                        fig_case_trend_total.update_layout(hovermode="x unified")


                        # --- Generate Heatmap (for PPT) ---
                        df_trend_location_actuals = df_case_details.groupby(['visit_month', 'location']).agg(
                            case_count=pd.NamedAgg(column='case_number', aggfunc='nunique')
                        ).reset_index()

                        # --- FIX: Create a complete data scaffold ---
                        all_months = pd.period_range(
                            start=df_case_details['visit_month'].min(),
                            end=df_case_details['visit_month'].max(),
                            freq='M'
                        )
                        all_locations_in_data = df_case_details['location'].unique()

                        new_index = pd.MultiIndex.from_product(
                            [all_months, all_locations_in_data],
                            names=['visit_month', 'location']
                        )
                        df_trend_complete = pd.DataFrame(index=new_index).reset_index()

                        df_trend_complete = pd.merge(
                            df_trend_complete,
                            df_trend_location_actuals,
                            on=['visit_month', 'location'],
                            how='left'
                        )

                        df_trend_complete['case_count'] = df_trend_complete['case_count'].fillna(0)
                        # --- END FIX ---

                        # --- FIX: Pivot the data for the heatmap ---
                        df_heatmap_pivot = df_trend_complete.pivot_table(
                            index='location',
                            columns='visit_month', # Use the period object for correct sorting
                            values='case_count',
                            fill_value=0 # Ensure all cells have a value
                        )

                        # Format the column names nicely (e.g., "2025-01")
                        df_heatmap_pivot.columns = df_heatmap_pivot.columns.to_timestamp().strftime('%Y-%m')

                        # Assign to fig_case_heatmap for PowerPoint
                        fig_case_heatmap = px.imshow(
                            df_heatmap_pivot, # <-- Pass the PIVOTED data
                            title='New Cases Heatmap by Location',
                            labels={'color': 'Number of Cases', 'x': 'Month', 'y': 'Location'}, # 'z' is now 'color'
                            color_continuous_scale=px.colors.sequential.Blues, # <-- NEW BLUE PALETTE
                            text_auto=True, # <-- Show the numbers
                            aspect="auto" # Adjust aspect ratio to fit container
                        )

                        fig_case_heatmap.update_xaxes(type='category', tickangle=-45) # Angle ticks for better fit
                        fig_case_heatmap.update_layout(
                             xaxis=dict(tickfont=dict(size=10)), # Smaller font for x-axis
                             yaxis=dict(tickfont=dict(size=10))  # Smaller font for y-axis
                        )
                        # --- End Heatmap Generation ---


                        # --- Display the selected chart in Streamlit ---
                        if trend_view == "Total Cases":
                            # FIX: use_container_width=True -> width='stretch'
                            st.plotly_chart(fig_case_trend_total, use_container_width=True)
                        else: # "Cases by Location"
                            # FIX: use_container_width=True -> width='stretch'
                            st.plotly_chart(fig_case_heatmap, use_container_width=True)

                    st.divider() # Add a divider after the graph

                    # 6. --- TABLES THIRD ---
                    case_data_col1, case_data_col2 = st.columns(2)
                    cost_label_suffix = f"({parts_label_suffix})"

                    with case_data_col1:
                        # --- UI/UX IMPROVEMENT: Added container ---
                        with st.container(border=True):
                            st.subheader("Top 10 Most Expensive Cases")
                            # FIX: use_container_width=True -> width='stretch'
                            st.dataframe(
                                df_case_agg.nlargest(10, 'total_cost_per_case'),
                                column_config={
                                    "total_cost_per_case": st.column_config.NumberColumn(format="$%.2f"),
                                    "parts_cost_per_case": st.column_config.NumberColumn(f"Parts Cost {cost_label_suffix}", format="$%.2f"),
                                    "labor_cost_per_case": st.column_config.NumberColumn(f"Labor Cost {cost_label_suffix}", format="$%.2f"),
                                    "total_hours_per_case": st.column_config.NumberColumn(format="%.1f h")
                                },
                                use_container_width=True,
                                hide_index=True
                            )

                    with case_data_col2:
                        # --- UI/UX IMPROVEMENT: Added container ---
                        with st.container(border=True):
                            st.subheader("Top 10 Cases by Most Visits")
                            # FIX: use_container_width=True -> width='stretch'
                            st.dataframe(
                                df_case_agg.nlargest(10, 'visits_per_case'),
                                 column_config={
                                     "total_cost_per_case": st.column_config.NumberColumn(format="$%.2f"),
                                     "parts_cost_per_case": st.column_config.NumberColumn(f"Parts Cost {cost_label_suffix}", format="$%.2f"),
                                     "labor_cost_per_case": st.column_config.NumberColumn(f"Labor Cost {cost_label_suffix}", format="$%.2f"),
                                     "total_hours_per_case": st.column_config.NumberColumn(format="%.1f h")
                                 },
                                 use_container_width=True,
                                 hide_index=True
                            )

            # 7. Raw Data
            with tab_data:
                st.title("💾 Raw Data Views")

                if st.checkbox('Show Raw Parts Data Only'):
                    st.subheader("Raw Data: Parts Lines (Filtered)")
                    st.dataframe(parts_df)

                if st.checkbox('Show Raw Labor Data Only'):
                    st.subheader("Raw Data: Labor Lines (Filtered)")
                    st.dataframe(labor_df)

                if st.checkbox('Show All Raw Service Data (Labor & Parts)'):
                    st.subheader("Raw Data: All Lines (Filtered)")
                    st.dataframe(full_df)

                st.caption(f"Showing **{full_df.shape[0]}** lines after all filters have been applied.")

            # 8. --- NEW: AI NARRATIVE TAB ---
            with tab_ai_narrative:
                st.title("💡 AI-Powered Value Story")
                st.markdown(
                    """
                    This panel uses a locally running **Ollama** model (requires `ollama serve` and the specified model pulled)
                    to analyze your filtered KPIs. It generates a tailored, persuasive narrative demonstrating the value 
                    of a service agreement to key stakeholders (Administrators, Oncologists, and Physicists).
                    """
                )
                st.divider()

                # Allow user to select their running model
                model_name = st.text_input(
                    "Enter your Ollama Model Name:",
                    value="gemma3",  # <-- UPDATED DEFAULT MODEL
                    help="Ensure this model is running in your local Ollama instance (e.g., 'gemma3', 'llama3', 'mistral')."
                )

                if st.button("Generate AI Narrative", type="primary", use_container_width=True):
                    if full_df_filtered.empty:
                        st.error("No data to analyze. Please adjust filters or upload a valid file.")
                    else:
                        with st.spinner(f"Connecting to Ollama and generating insights with '{model_name}'..."):
                            try:
                                # 1. Aggregate Top Parts Data (use gross cost for "sticker shock" value)
                                # Only do this if parts_df is not empty
                                if not parts_df.empty:
                                    top_parts_data = parts_df.groupby('item').agg(
                                        total_cost=pd.NamedAgg(column='parts_gross_cost', aggfunc='sum'),
                                        total_qty=pd.NamedAgg(column='qty', aggfunc='sum')
                                    ).nlargest(5, 'total_cost').reset_index()
                                    # Rename columns for PPT input consistency
                                    top_parts_data.rename(columns={'total_cost': 'gross_cost', 'total_qty': 'qty'}, inplace=True)
                                else:
                                    top_parts_data = pd.DataFrame(columns=['item', 'gross_cost', 'qty'])


                                # 2. Aggregate Top Cases Data (using pre-calculated df_case_agg)
                                # Only do this if df_case_agg is not empty
                                if not df_case_agg.empty:
                                    top_cases_data = df_case_agg.nlargest(5, 'total_cost_per_case')[[
                                        'case_number', 'total_cost_per_case', 'visits_per_case', 'total_hours_per_case'
                                    ]]
                                else:
                                     top_cases_data = pd.DataFrame(columns=['case_number', 'total_cost_per_case', 'visits_per_case', 'total_hours_per_case'])

                                # 3. Aggregate Top Issues Data (for AI prompt - based on df_issue_summary)
                                top_issues_data_for_ai = df_issue_summary.head(5).copy()

                                # 4. Consolidate all KPIs into a dictionary
                                kpi_data_for_ai = {
                                    "period_start_date": str(start_date_input),
                                    "period_end_date": str(end_date_input),
                                    "cost_basis_note": f"Costs are calculated as {parts_label_suffix} Cost.",
                                    "tcs_label": tcs_label,
                                    "total_cost": total_tcs,
                                    "total_labor_cost": total_labor_current_cost,
                                    "total_parts_cost": total_parts_current_cost,
                                    "total_service_events": total_events,
                                    "average_cost_per_event": avg_tcs_per_event,
                                    "average_labor_cost_per_event": avg_labor_only_per_event, # <-- NEW KPI FOR AI
                                    "total_labor_hours": total_hours,
                                    "total_parts_replaced_qty": total_parts_replaced,
                                    "total_discount_given": total_discount_given,
                                    "total_reactive_cases": total_cases,
                                    "average_cost_per_case": avg_cost_per_case,
                                    "average_visits_per_case": avg_visits_per_case,
                                    "top_5_expensive_parts": top_parts_data.to_dict('records'),
                                    "top_5_expensive_cases": top_cases_data.to_dict('records'),
                                    "top_5_recurring_issues": top_issues_data_for_ai.to_dict('records') # <-- NEW KPI FOR AI
                                }

                                # 5. Call the AI function
                                stream = get_ai_narrative(kpi_data_for_ai, model_name)

                                # 6. Stream the response
                                st.subheader(f"AI-Generated Narrative (using {model_name})")
                                with st.container(border=True):
                                    # st.write_stream displays the streaming output
                                    def stream_handler():
                                        for chunk in stream:
                                            # Check if the chunk is from the chat response
                                            if 'message' in chunk and 'content' in chunk['message']:
                                                yield chunk['message']['content']
                                            # Check if it's the custom error message
                                            elif isinstance(chunk, str):
                                                yield chunk

                                    st.write_stream(stream_handler)

                            except Exception as e:
                                st.error(f"An error occurred while communicating with Ollama: {e}")
                                import traceback
                                st.code(traceback.format_exc())
            # --- END: AI NARRATIVE TAB ---


            # --- ADD POWERPOINT DOWNLOAD BUTTON (Must be INSIDE the main try block) ---
            st.sidebar.divider()
            st.sidebar.subheader("Download Full Report")

            if st.sidebar.button("Generate PowerPoint Report"):
                with st.spinner("Generating PowerPoint... This may take a moment."):
                    # Create the date range string for the title slide
                    date_str = (
                        f"{START_DATE_FILTER.strftime('%Y-%m-%d')} to "
                        f"{(END_DATE_FILTER - timedelta(days=1)).strftime('%Y-%m-%d')}"
                    )
                    
                    # --- Aggregate Final Data for PPT ---
                    
                    # 1. Top 5 Most Expensive Parts (for Parts Slide)
                    # Use the df_parts_cost calculated in tab_parts (ensure to calculate this before the PPT button is pressed)
                    if not parts_df.empty:
                        # Re-calculate to ensure it's fresh after filtering
                        df_parts_cost_final = parts_df.groupby('item').agg(
                            gross_cost=pd.NamedAgg(column='parts_gross_cost', aggfunc='sum'),
                            qty=pd.NamedAgg(column='qty', aggfunc='sum')
                        ).sort_values(by='gross_cost', ascending=False).head(5).reset_index()
                        summary_top_parts_for_ppt = df_parts_cost_final.to_dict('records')
                    else:
                        summary_top_parts_for_ppt = []

                    # 2. Top 5 Most Frequent Corrective Actions (for Activity Slide)
                    summary_top_issues_for_ppt = df_issue_summary.head(5).to_dict('records')
                    
                    # --- END Aggregate Final Data for PPT ---


                    # Call the generation function with all your figures and KPIs
                    ppt_data = generate_powerpoint_report(
                        # Figures
                        fig_kpi_trend=fig_trend, # from tab_kpi
                        fig_cost_split=fig_split, # from tab_kpi
                        fig_tech=fig_tech, # from tab_performance
                        fig_loc=fig_loc, # from tab_performance
                        fig_activity=fig_activity, # from tab_activity
                        fig_parts_qty=fig_parts_qty, # from tab_parts
                        fig_parts_cost=fig_parts_cost, # from tab_parts
                        fig_case_trend_total=fig_case_trend_total, # from tab_case
                        fig_case_heatmap=fig_case_heatmap, # from tab_case
                        # Report Details
                        report_title="Labor & Service Analysis Report",
                        date_range_str=date_str,
                        # Main KPIs
                        kpi_total_tcs=total_tcs,
                        kpi_tcs_label=tcs_label,
                        kpi_labor_cost=total_labor_current_cost,
                        kpi_parts_cost=total_parts_current_cost,
                        kpi_labor_label=f"Labor Cost ({parts_label_suffix})",
                        kpi_parts_label=parts_label_suffix,
                        kpi_total_events=total_events,
                        kpi_avg_tcs=avg_tcs_per_event,
                        kpi_total_hours=total_hours,
                        kpi_total_parts=total_parts_replaced,
                        # Case KPIs
                        kpi_total_cases=total_cases,
                        kpi_avg_cost_case=avg_cost_per_case,
                        kpi_avg_visits_case=avg_visits_per_case,
                        # NEW TALKING POINT DATA
                        kpi_total_discount=total_discount_given, 
                        summary_top_parts=summary_top_parts_for_ppt,
                        summary_top_issues=summary_top_issues_for_ppt
                    )

                # Provide the download button
                st.sidebar.download_button(
                    label="Download .pptx Report",
                    data=ppt_data,
                    file_name="Service_Analysis_Report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key='ppt_download' # Add a key for stability
                )
                st.sidebar.success("Report is ready! Click 'Download .pptx Report'.")


        except Exception as e:
            st.error(f"An error occurred during data processing or report generation: {e}")
            import traceback
            st.code(f"Error Details: {e}\n{traceback.format_exc()}", language='text')