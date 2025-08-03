import streamlit as st
import pandas as pd
import plotly.express as px
from datetime import datetime, timedelta
import os
import random
from io import BytesIO

# Import pptx libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN # Import PP_ALIGN for text alignment

# --- Define Color Scheme ---
PRIMARY_COLOR = 'rgb(43, 101, 124)'
SECONDARY_COLOR = 'rgb(54, 164, 179)'
ELEKTA_FONT_COLOR = RGBColor(43,101,125) # From user's pptx code
RGB_CUSTOM_COLORS = [ # From user's pptx code
    'rgb(43,101,125)',  # Base color (teal-blue)
    'rgb(85,130,145)',  # Lighter and more saturated
    'rgb(25,85,105)',   # Darker and less saturated
    'rgb(135,175,195)', # Much lighter
    'rgb(0,60,80)',     # Very dark teal
    'rgb(58,121,150)',  # Brighter, with more blue
    'rgb(90,140,160)',  # Muted teal-blue
    'rgb(10,70,90)',    # Dark, almost navy
    'rgb(100,180,200)', # Light cyan-blue
    'rgb(30,90,110)'    # Deep teal
]
COLOR_SEQUENCE = RGB_CUSTOM_COLORS # Use this for Plotly charts for consistency

# Set title and configure Streamlit page layout
TITLE = 'Service Agreement Report'
DOWNTIME_URL = 'https://elekta.lightning.force.com/lightning/r/Report/00OKf000000Z339MAC/view?queryScope=userFolders'

# --- Helper functions for PPTX ---
def add_custom_textbox(slide, left:Inches, top:Inches, width: Inches, height: Inches, font_name: str, font_size:Pt, font_color: RGBColor, bold: bool, text: str, text_align: PP_ALIGN = PP_ALIGN.LEFT):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.name = font_name
    text_frame.paragraphs[0].font.size = font_size
    text_frame.paragraphs[0].font.bold = bold
    text_frame.paragraphs[0].font.color.rgb = font_color
    text_frame.paragraphs[0].alignment = text_align # Set text alignment

def add_formatted_text_line(slide, left:Inches, top:Inches, width: Inches, height: Inches, font_name: str, font_size:Pt, font_color: RGBColor, text_parts: list[tuple[str, bool]], text_align: PP_ALIGN = PP_ALIGN.LEFT):
    """
    Adds a textbox with a single paragraph that can contain multiple runs with different bolding.
    text_parts: A list of tuples, where each tuple is (text_string, is_bold_boolean).
    """
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.alignment = text_align

    for text, bold_status in text_parts:
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = font_name
        font.size = font_size
        font.bold = bold_status
        font.color.rgb = font_color

def add_rectangle_background(slide, left:Inches, top:Inches, width:Inches, height:Inches, BGcolor:RGBColor, border:int):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = BGcolor
    shape.shadow.inherit = False
    shape.line.fill.background()  # Inherit shadow from the slide
    
    # # --- MODIFIED: Robustly handle line properties based on 'border' parameter ---
    # line = shape.line # Get the line object
    # if border == 0:
    #     line.fill.background() # No line
    # else:
    #     line.fill.solid() # Solid line fill
    #     line.fore_color.rgb = RGBColor(0, 0, 0) # Black border
    #     line.width = Pt(border) # Border width (Pt() handles integer or float correctly)
    # --- END MODIFIED ---

def add_table_to_slide(slide, df_table, left, top, width, height, font_name, font_size):
    rows, cols = df_table.shape
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table # +1 for header row

    # Set column widths (optional, but helps layout)
    for i, col_width in enumerate([width / cols] * cols):
        table.columns[i].width = col_width

    # Set header row
    for col_idx, col_name in enumerate(df_table.columns):
        cell = table.cell(0, col_idx)
        text_frame = cell.text_frame
        text_frame.text = str(col_name)
        text_frame.paragraphs[0].font.name = font_name
        text_frame.paragraphs[0].font.size = Pt(font_size)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(230, 230, 230) # Light grey header

    # Populate data rows
    for row_idx, row_data in df_table.iterrows():
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            text_frame = cell.text_frame
            text_frame.text = str(cell_value)
            text_frame.paragraphs[0].font.name = font_name
            text_frame.paragraphs[0].font.size = Pt(font_size)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def generate_service_contract_slides(df_data: pd.DataFrame, ppt_title: str):
    prs = Presentation()
    prs.slide_width = Inches(26.66)
    prs.slide_height = Inches(15)
    font_name = 'Calibri'
    image_folder_ = './images/'
    image_folder = './images/Cards'
    output_graph_dir = 'graphs/service_dashboard'
    os.makedirs(output_graph_dir, exist_ok=True)
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    slide_layout = prs.slide_layouts[6]

    # --- 1. Title Slide ---
    slide = prs.slides.add_slide(slide_layout)
    try:
        if os.path.exists(image_folder_) and os.listdir(image_folder_):
            images_in_folder = [f for f in os.listdir(image_folder_) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))]
            if images_in_folder:
                random_image = random.choice(images_in_folder)
                image_path = os.path.join(image_folder_, random_image)
                slide.shapes.add_picture(image_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
            else:
                st.warning("No image files found in the 'images' folder for the title slide background.")
        else:
            st.warning("The 'images' folder does not exist or is empty for the title slide background.")
    except Exception as e:
        st.error(f"Error adding background image to title slide: {e}")

    add_custom_textbox(slide=slide, left=Inches(1), top=Inches(5.47), width=Inches(11), height=Inches(3),
                       font_name=font_name, font_color=ELEKTA_FONT_COLOR, font_size=Pt(90), bold=True, text=ppt_title)

    # --- 2. KPI Summary Slide ---
    slide = prs.slides.add_slide(slide_layout)
    add_rectangle_background(slide, Inches(0.81), Inches(2.7), Inches(24.75), Inches(11.86), RGBColor(248,248,248), 0)
    add_custom_textbox(slide, Inches(0.8), Inches(1.08), Inches(24), Inches(2.9), font_name, Pt(80), ELEKTA_FONT_COLOR, True, "Service Contract Key Metrics")

    total_devices = df_data.shape[0]
    contracts_expiring_soon = df_data[df_data['Contract Status'] == 'Expiring Soon'].shape[0]
    expired_contracts = df_data[df_data['Contract Status'] == 'Expired'].shape[0]
    total_contract_value = df_data['Contract Price'].sum() if 'Contract Price' in df_data.columns else 0

    add_custom_textbox(slide, Inches(2), Inches(4), Inches(5), Inches(1), font_name, Pt(40), ELEKTA_FONT_COLOR, False, "Total Devices Installed")
    add_custom_textbox(slide, Inches(3), Inches(4.7), Inches(5), Inches(1), font_name, Pt(100), ELEKTA_FONT_COLOR, True, f"{total_devices}")

    add_custom_textbox(slide, Inches(9), Inches(4), Inches(5), Inches(1), font_name, Pt(40), ELEKTA_FONT_COLOR, False, "Contracts Expiring (<12 weeks)")
    add_custom_textbox(slide, Inches(10), Inches(4.7), Inches(5), Inches(1), font_name, Pt(100), ELEKTA_FONT_COLOR, True, f"{contracts_expiring_soon}")

    add_custom_textbox(slide, Inches(16), Inches(4), Inches(5), Inches(1), font_name, Pt(40), ELEKTA_FONT_COLOR, False, "Expired Contracts")
    add_custom_textbox(slide, Inches(17), Inches(4.7), Inches(5), Inches(1), font_name, Pt(100), ELEKTA_FONT_COLOR, True, f"{expired_contracts}")

    add_custom_textbox(slide, Inches(2), Inches(9), Inches(5), Inches(1), font_name, Pt(60), ELEKTA_FONT_COLOR, False, "Total Contract Annual Value")
    add_custom_textbox(slide, Inches(2), Inches(10.2), Inches(5), Inches(1), font_name, Pt(120), ELEKTA_FONT_COLOR, True, f"${total_contract_value:,.0f}")

    # --- 3. Contract Status Distribution Slide ---
    slide = prs.slides.add_slide(slide_layout)
    add_rectangle_background(slide, Inches(0.81), Inches(2.7), Inches(24.75), Inches(11.86), RGBColor(248,248,248), 0)
    add_custom_textbox(slide, Inches(1.27), Inches(1.08), Inches(24), Inches(2.9), font_name, Pt(70), ELEKTA_FONT_COLOR, True, "Contract Status Distribution")
    
    contract_status_counts = df_data['Contract Status'].value_counts().reset_index()
    contract_status_counts.columns = ['Status', 'Count']
    fig_contract_status = px.pie(contract_status_counts, values='Count', names='Status',
                                 title='Overall Contract Status Distribution',
                                 color_discrete_sequence=COLOR_SEQUENCE,
                                 template="plotly_white")
    fig_contract_status.update_layout(width=1200, height=800)
    contract_status_path = os.path.join(output_graph_dir, f'contract_status_pie_{timestamp}.png')
    fig_contract_status.write_image(contract_status_path)
    slide.shapes.add_picture(contract_status_path, Inches(5), Inches(3.5), width=Inches(16))

    # --- 4. Device Lifecycle & Risk (Age Distribution) Slide ---
    if 'Device Age Group' in df_data.columns:
        slide = prs.slides.add_slide(slide_layout)
        add_rectangle_background(slide, Inches(0.81), Inches(2.7), Inches(24.75), Inches(11.86), RGBColor(248,248,248), 0)
        add_custom_textbox(slide, Inches(1.27), Inches(1.08), Inches(24), Inches(2.9), font_name, Pt(70), ELEKTA_FONT_COLOR, True, "Device Lifecycle & Risk")
        
        fig_age = px.histogram(df_data, x='Device Age Group',
                               title='Distribution of Device Ages',
                               labels={'Device Age Group': 'Device Age (Years)'},
                               category_orders={"Device Age Group": ["0-5 years", "5-10 years", ">10 years"]},
                               color_discrete_sequence=COLOR_SEQUENCE,
                               template="plotly_white")
        fig_age.update_layout(bargap=0.8, showlegend=True)
        fig_age.update_layout(width=1800, height=900)
        
        age_hist_path = os.path.join(output_graph_dir, f'device_age_histogram_{timestamp}.png')
        fig_age.write_image(age_hist_path)
        slide.shapes.add_picture(age_hist_path, Inches(2), Inches(4), width=Inches(20))

    # --- 5. Upcoming Renewals & Expirations Slide ---
    if 'Weeks To Renewal' in df_data.columns:
        slide = prs.slides.add_slide(slide_layout)
        add_rectangle_background(slide, Inches(0.81), Inches(2.7), Inches(24.75), Inches(11.86), RGBColor(248,248,248), 0)
        add_custom_textbox(slide, Inches(0.8), Inches(1.08), Inches(24), Inches(2.9), font_name, Pt(80), ELEKTA_FONT_COLOR, True, "Upcoming Renewals & Expirations")

        upcoming_renewals = df_data[df_data['Weeks To Renewal'] <= 52].copy() # Filter for contracts expiring in the next 52 weeks
        if not upcoming_renewals.empty:
            # Create bins for weeks to renewal, ensuring monotonicity
            upcoming_renewals['Renewal Period'] = pd.cut(upcoming_renewals['Weeks To Renewal'],
                                                         bins=[-0.1, 0, 12, 26, 52], # Adjusted to capture 0 correctly
                                                         labels=['Expired', '0-12 Weeks', '12-26 Weeks', '26-52 Weeks'],
                                                         right=True, # (a, b]
                                                         include_lowest=True)
            # Add a category for values outside these bins (e.g., >52 weeks if not filtered)
            upcoming_renewals['Renewal Period'] = upcoming_renewals['Renewal Period'].cat.add_categories('>52 Weeks').fillna('>52 Weeks')
            upcoming_renewals = upcoming_renewals[upcoming_renewals['Renewal Period'] != '>52 Weeks'] # Filter out if not needed for chart
            
            # Group by 'Location' instead of 'Account'
            renewal_counts = upcoming_renewals.groupby(['Renewal Period','Location']).size().unstack(fill_value=0)
            
            fig_renewals = px.bar(renewal_counts, x=renewal_counts.index, y=renewal_counts.columns,
                                  title='Number of Contracts by Upcoming Renewal Period',
                                  labels={'value': 'Number of Contracts', 'Location': 'Location'},
                                  color_discrete_sequence=COLOR_SEQUENCE,
                                  template="plotly_white")
            fig_renewals.update_layout(barmode='stack', width=1800, height=900)
            
            renewals_path = os.path.join(output_graph_dir, f'upcoming_renewals_bar_{timestamp}.png')
            fig_renewals.write_image(renewals_path)
            slide.shapes.add_picture(renewals_path, Inches(2), Inches(4), width=Inches(20))
        else:
            add_custom_textbox(slide, Inches(2), Inches(5), Inches(20), Inches(2), font_name, Pt(30), RGBColor(100,100,100), False, "No upcoming renewals in the next 52 weeks.", text_align=PP_ALIGN.CENTER)

    # --- 6. Financial Summary Slide ---
    if 'Contract Price' in df_data.columns:
        slide = prs.slides.add_slide(slide_layout)
        add_rectangle_background(slide, Inches(0.81), Inches(2.7), Inches(24.75), Inches(11.86), RGBColor(248,248,248), 0)
        add_custom_textbox(slide, Inches(1.27), Inches(1.08), Inches(24), Inches(2.9), font_name, Pt(70), ELEKTA_FONT_COLOR, True, "Contract Value by Location")

        # Group by 'Location' instead of 'Account'
        contract_value_by_location = df_data.groupby('Location')['Contract Price'].sum().reset_index()
        contract_value_by_location = contract_value_by_location.sort_values(by='Contract Price', ascending=False)
        
        fig_financial = px.bar(contract_value_by_location, y='Location', x='Contract Price',
                               title='Total Contract Value by Location',
                               labels={'Contract Price': 'Contract Value'},
                               color='Location', # Color by 'Location'
                               color_discrete_sequence=COLOR_SEQUENCE,
                               template="plotly_white")
        fig_financial.update_layout(width=1800, height=900, showlegend=False)
        
        financial_path = os.path.join(output_graph_dir, f'financial_summary_bar_{timestamp}.png')
        fig_financial.write_image(financial_path)
        slide.shapes.add_picture(financial_path, Inches(2), Inches(4), width=Inches(20))

    # # --- 7. Devices at Risk Table Slide ---
    # slide = prs.slides.add_slide(slide_layout)
    # add_rectangle_background(slide, Inches(0.81), Inches(2.7), Inches(24.75), Inches(11.86), RGBColor(248,248,248), 0)
    # add_custom_textbox(slide, Inches(1.27), Inches(1.08), Inches(24), Inches(2.9), font_name, Pt(70), ELEKTA_FONT_COLOR, True, "Devices At Risk")

    # at_risk_devices = df_data[(df_data['Contract Status'] == 'Expired') | (df_data['Contract Status'] == 'Expiring Soon')].copy()
    # if not at_risk_devices.empty:
    #     # Select relevant columns for the table
    #     table_df = at_risk_devices[['Account', 'Location', 'Display Product Name', 'Serial Number', 'Contract Status', 'Weeks To Renewal', 'Warranty End Date', 'EoL Date IP']]
    #     # Convert dates to string format for display
    #     table_df['Warranty End Date'] = table_df['Warranty End Date'].dt.strftime('%m/%d/%Y').fillna('N/A')
    #     table_df['EoL Date IP'] = table_df['EoL Date IP'].dt.strftime('%m/%d/%Y').fillna('N/A')

    #     add_table_to_slide(slide, table_df, Inches(1.5), Inches(4), Inches(23), Inches(9), font_name, Pt(10))
    # else:
    #     add_custom_textbox(slide, Inches(2), Inches(5), Inches(20), Inches(2), font_name, Pt(30), RGBColor(100,100,100), False, "No devices currently at risk.", text_align=PP_ALIGN.CENTER)

    # --- 8. Individual Device Cards Slides for PowerPoint ---
    # Define card dimensions and positions for 3 cards per slide
    card_width_ppt = Inches(7.5) # Adjusted width for each card in PPT
    card_height_ppt = Inches(9) # Height of each card area in PPT
    
    # Horizontal starting positions for each of the 3 cards on a slide
    card_starts_x_ppt = [Inches(1), Inches(9.5), Inches(18)]
    
    # Vertical starting position for cards on a slide
    card_start_y_ppt = Inches(4) 

    # Loop through devices in chunks of 3
    for i in range(0, len(df_data), 3):
        slide = prs.slides.add_slide(slide_layout)
        add_rectangle_background(slide, Inches(0.3), Inches(2.7), Inches(26.00), Inches(11.86), RGBColor(248,248,248), 0) # Background for the whole slide

        add_custom_textbox(slide, Inches(0.8), Inches(1.08), Inches(24), Inches(1.5), font_name, Pt(80), ELEKTA_FONT_COLOR, True, "Machine Fleet Overview")

        devices_on_this_slide = df_data.iloc[i : i + 3]

        for j, (idx, device) in enumerate(devices_on_this_slide.iterrows()):
            current_card_left_ppt = card_starts_x_ppt[j]
            
            # Draw a rectangle for the card background
            add_rectangle_background(slide, current_card_left_ppt, card_start_y_ppt, card_width_ppt, card_height_ppt, RGBColor(255,255,255),0) # White background, with border

            # Prepare Customs Acceptance Date string
            customs_acceptance_date_str = "N/A"
            if 'Customs Acceptance Date' in device.index and pd.notna(device.get('Customs Acceptance Date')):
                customs_acceptance_date_str = device['Customs Acceptance Date'].strftime('%m/%d/%Y')
            
            # Prepare Contract End Date string for "Contract Expires"
            end_date_str = "N/A"
            if 'Contract End Date' in device.index and pd.notna(device.get('Contract End Date')):
                end_date_str = device['Contract End Date'].strftime('%m/%d/%Y')


            original_product_string = device.get('Installed Product', 'Unknown')
            first_part_for_image = original_product_string.split('/')[0] if isinstance(original_product_string, str) else 'Unknown'
            sanitized_filename_part = "".join(c if c.isalnum() else '_' for c in first_part_for_image).replace('__', '_').strip('_')
            image_filename = f"{sanitized_filename_part}.png"

            full_image_path_on_disk = os.path.join(image_folder, image_filename)
            
            # Add and center device image within its card
            img_width_card_ppt = Inches(3) # Width of the image within the card
            img_left_card_ppt = current_card_left_ppt + (card_width_ppt - img_width_card_ppt) / 2 # Centered within the card's width
            img_top_card_ppt = card_start_y_ppt + Inches(0.5) # A little padding from the top of the card
            
            if os.path.exists(full_image_path_on_disk):
                slide.shapes.add_picture(full_image_path_on_disk, img_left_card_ppt, img_top_card_ppt, width=img_width_card_ppt)
            else:
                # Placeholder for PPTX if image not found
                # Create a simple text placeholder instead of a complex image URL
                add_custom_textbox(slide, img_left_card_ppt, img_top_card_ppt, img_width_card_ppt, Inches(1),
                                   font_name, Pt(10), RGBColor(150,150,150), False, "Image N/A", text_align=PP_ALIGN.CENTER)
                st.warning(f"PPTX Image not found for '{first_part_for_image}'. Looked for: {full_image_path_on_disk}")

            # Add device details (text boxes) within its card
            text_box_left_padding_ppt = Inches(0.5) # Padding from the left edge of the card
            text_box_width_ppt = card_width_ppt - (2 * text_box_left_padding_ppt) # Text box spans most of the card width

            # Position for the product name (title)
            add_custom_textbox(slide,
                               left=current_card_left_ppt + text_box_left_padding_ppt,
                               top=img_top_card_ppt + img_width_card_ppt + Inches(0.2), # Below the image
                               width=text_box_width_ppt,
                               height=Inches(0.8),
                               font_name=font_name,
                               font_size=Pt(40), # Smaller font for card text
                               font_color=ELEKTA_FONT_COLOR,
                               bold=True,
                               text=f"{device.get('Display Product Name', 'N/A')}",
                               text_align=PP_ALIGN.CENTER) # Center the product name

            # Positions for other details
            detail_start_top_ppt = img_top_card_ppt + img_width_card_ppt + Inches(1.2) # Start below the title
            line_height_ppt = Inches(0.5) # Space between lines of text

            add_formatted_text_line(slide,
                               left=current_card_left_ppt + text_box_left_padding_ppt,
                               top=detail_start_top_ppt,
                               width=text_box_width_ppt,
                               height=line_height_ppt,
                               font_name=font_name,
                               font_size=Pt(25), # Consistent font size
                               font_color=RGBColor(50,50,50),
                               text_parts=[("Contract Expires: ", True), (end_date_str, False)],
                               text_align=PP_ALIGN.CENTER) # Center customs acceptance date

            add_formatted_text_line(slide,
                               left=current_card_left_ppt + text_box_left_padding_ppt,
                               top=detail_start_top_ppt + line_height_ppt,
                               width=text_box_width_ppt,
                               height=line_height_ppt,
                               font_name=font_name,
                               font_size=Pt(25), # Consistent font size
                               font_color=RGBColor(50,50,50),
                               text_parts=[("Age: ", True), (f"{device.get('Device Age', 'N/A')} years", False)],
                               text_align=PP_ALIGN.CENTER) # Center device age

            add_formatted_text_line(slide,
                               left=current_card_left_ppt + text_box_left_padding_ppt,
                               top=detail_start_top_ppt + 2*line_height_ppt,
                               width=text_box_width_ppt,
                               height=line_height_ppt,
                               font_name=font_name,
                               font_size=Pt(25), # Consistent font size
                               font_color=RGBColor(50,50,50),
                               text_parts=[("Renew In: ", True), (f"{device.get('Weeks To Renewal', 'N/A')} weeks", False)],
                               text_align=PP_ALIGN.CENTER) # Center customs acceptance date
            
            # Add a horizontal line (hr equivalent)
            line_top_ppt = detail_start_top_ppt + 3 * line_height_ppt + Inches(0.4)
            line_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, current_card_left_ppt + text_box_left_padding_ppt, line_top_ppt,
                text_box_width_ppt, Inches(0.02) # Very thin line
            )
            # line_shape.fill.solid().fore_color.rgb = RGBColor(200, 200, 200) # Ensure line color is set

            add_formatted_text_line(slide,
                               left=current_card_left_ppt + text_box_left_padding_ppt,
                               top=detail_start_top_ppt + 4*line_height_ppt + Inches(0.3), # Below the line
                               width=text_box_width_ppt,
                               height=line_height_ppt,
                               font_name=font_name,
                               font_size=Pt(20), # Consistent font size
                               font_color=RGBColor(100,100,100),
                               text_parts=[("CAT: ", True), (customs_acceptance_date_str, False)],
                               text_align=PP_ALIGN.CENTER) # Center customs acceptance date
            
            add_formatted_text_line(slide,
                               left=current_card_left_ppt + text_box_left_padding_ppt,
                               top=detail_start_top_ppt + 5 * line_height_ppt + Inches(0.3), # Below the line
                               width=text_box_width_ppt,
                               height=line_height_ppt,
                               font_name=font_name,
                               font_size=Pt(20), # Consistent font size
                               font_color=RGBColor(100,100,100),
                               text_parts=[("Warranty End Date: ", True), (device.get('Warranty End Date', 'N/A').strftime('%m/%d/%Y') if pd.notna(device.get('Warranty End Date')) else 'N/A', False)],
                               text_align=PP_ALIGN.CENTER) # Center warranty end date

            add_formatted_text_line(slide,
                               left=current_card_left_ppt + text_box_left_padding_ppt,
                               top=detail_start_top_ppt + 6 * line_height_ppt + Inches(0.3),
                               width=text_box_width_ppt,
                               height=line_height_ppt,
                               font_name=font_name,
                               font_size=Pt(20), # Consistent font size
                               font_color=RGBColor(100,100,100),
                               text_parts=[("EoL Date IP: ", True), (device.get('EoL Date IP', 'N/A').strftime('%m/%d/%Y') if pd.notna(device.get('EoL Date IP')) else 'N/A', False)],
                               text_align=PP_ALIGN.CENTER) # Center EoL date

    # Save the presentation
    output_filename = f'{ppt_title.replace(" ", "_")}_{timestamp}.pptx'
    prs.save(f"presentations/Service_Agreements/{output_filename}")
    return output_filename


# --- Page Configuration ---
st.set_page_config(
    page_title="Service Contracts Dashboard",
    page_icon="🏥",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- Custom CSS for a cleaner look and matching color scheme ---
st.markdown(f"""
    <style>
    .main {{
        background-color: #f0f2f6;
        padding: 20px;
    }}
    .st-emotion-cache-z5fcl4 {{ /* Adjust padding for main content area */
        padding-top: 2rem;
    }}
    .st-emotion-cache-1cyp85f {{ /* Adjust padding for columns */
        padding-top: 0rem;
        padding-bottom: 0rem;
    }}
    .st-emotion-cache-183p0q {{ /* Adjust font size for headers */
        font-size: 1.2rem;
        font-weight: bold;
        color: #262730;
    }}
    .block-container {{ /* General block container padding */
        padding-top: 1rem;
        padding-bottom: 0rem;
        padding-left: 1rem;
        padding-right: 1rem;
    }}
    .stButton>button {{ /* Styling for buttons */
        background-color: {PRIMARY_COLOR}; /* Match primary color */
        color: white;
        border-radius: 8px;
        padding: 10px 20px;
        font-size: 16px;
        border: none;
        cursor: pointer;
        transition: all 0.3s ease;
    }}
    .stButton>button:hover {{ /* Hover effect for buttons */
        background-color: {SECONDARY_COLOR}; /* Lighter hover color */
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.2);
    }}
    .stFileUploader {{ /* Styling for file uploader */
        border: 2px dashed {PRIMARY_COLOR}; /* Match primary color */
        border-radius: 8px;
        padding: 20px;
        text-align: center;
        background-color: #e6ffe6; /* Light green background */
    }}
    .stFileUploader label {{ /* Styling for file uploader label */
        color: {PRIMARY_COLOR}; /* Match primary color */
        font-weight: bold;
    }}

    /* --- CSS for Device Cards (now primarily handled by st.container) --- */
    /* The .device-card-content class is no longer explicitly used for the main card div,
       but the text-align: center is applied via inline style in st.markdown. */
    .st-emotion-cache-1r6dm1s {{ /* This targets the outer div of st.container(border=True) */
        background-color: #f8f8f8; /* Light grey background for cards */
        border-radius: 8px;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1); /* Subtle shadow */
        margin-bottom: 15px; /* Space between cards in a column */
        padding: 15px; /* Add padding to the container itself */
    }}
    .st-emotion-cache-1r6dm1s h3 {{ /* Target h3 inside the container */
        color: {PRIMARY_COLOR}; /* Make subheader match primary color */
        margin-top: 0;
        margin-bottom: 10px;
        font-size: 1.3em;
    }}
    .st-emotion-cache-1r6dm1s p {{ /* Target p inside the container */
        font-size: 0.95em;
        line-height: 1.4;
        margin-bottom: 5px;
    }}
    .st-emotion-cache-1r6dm1s strong {{ /* Target strong inside the container */
        color: #333; /* Darker text for emphasis */
    }}
    /* No specific .device-card-image needed as st.image handles its own styling/centering */
    </style>
""", unsafe_allow_html=True)

# --- Title and Description ---
st.markdown(f"<h1 style='color: {PRIMARY_COLOR};'>🏥 Service Contracts Dashboard</h1>", unsafe_allow_html=True)
st.markdown(
    """
    Upload your CSV/Excel file to visualize key insights about your radiation therapy equipment service contracts.
    This dashboard helps clinical staff and managers track contract status, device lifecycle, and financial impact
    for better operational and strategic planning.
    """
)

# --- File Uploader (in sidebar) ---
st.sidebar.title('Data Upload & Filters')

# Check if DataFrame is already in session state
if 'uploaded_df' not in st.session_state:
    uploaded_file = st.sidebar.file_uploader("Upload your data file", type=["csv", "xlsx", "xls"])
else:
    # If data is already uploaded, show a message and allow re-upload
    st.sidebar.success("Data already loaded. Upload a new file to replace.")
    uploaded_file = st.sidebar.file_uploader("Upload new data file", type=["csv", "xlsx", "xls"])

st.sidebar.divider()

df = None
if uploaded_file is not None:
    try:
        # Read the uploaded Excel or CSV file
        file_extension = uploaded_file.name.split('.')[-1]
        if file_extension in ['xlsx', 'xls']:
            df = pd.read_excel(uploaded_file)
        elif file_extension == 'csv':
            try:
                df = pd.read_csv(uploaded_file, encoding='utf-8')
            except UnicodeDecodeError:
                df = pd.read_csv(uploaded_file, encoding='ISO-8859-1')
        
        # Store the loaded DataFrame in session state
        st.session_state['uploaded_df'] = df.copy()
        st.sidebar.success("File uploaded successfully!")

    except Exception as e:
        st.error(f"Error processing file: {e}. Please ensure it's a valid CSV/Excel with expected columns.")
        st.info("Expected columns: 'Account', 'Location', 'Installed Product: Installed Product', 'Installed Product: Serial/Lot Number', 'Service/Maintenance Contract: Contract Name/Number', 'Installed Product: Customer/Device Acceptance Date', 'Covered Product: Record Number', 'Installed Product: Device Age', 'Current Term Start Date', 'Current Term End Date', 'Start Date', 'End Date', 'Weeks To Renewal', 'Installed Product: Warranty End Date', 'Installed Product: EoL Date IP', 'Installed Product: EoGS Date IP', 'Contract Status', 'SLA Terms', 'Contract Price Currency', 'Contract Price'")
        # Clear session state if file upload fails
        if 'uploaded_df' in st.session_state:
            del st.session_state['uploaded_df']

# Use the DataFrame from session state if available
if 'uploaded_df' in st.session_state and not st.session_state['uploaded_df'].empty:
    df = st.session_state['uploaded_df'].copy() # Work with a copy to avoid modifying session state directly

    # --- Data Preprocessing ---
    # Rename columns for consistency and easier access
    if 'Installed Product: Installed Product' in df.columns:
        df.rename(columns={'Installed Product: Installed Product': 'Installed Product'}, inplace=True)
    if 'Installed Product: Serial/Lot Number' in df.columns: # NEW: Rename this column
        df.rename(columns={'Installed Product: Serial/Lot Number': 'Serial Number'}, inplace=True)
    elif 'Serial/Lot Number' in df.columns: # Keep fallback for older column name
        df.rename(columns={'Serial/Lot Number': 'Serial Number'}, inplace=True)
    if 'Installed Product: Warranty End Date' in df.columns: # NEW: Rename this column
        df.rename(columns={'Installed Product: Warranty End Date': 'Warranty End Date'}, inplace=True)
    if 'Installed Product: EoL Date IP' in df.columns: # NEW: Rename this column
        df.rename(columns={'Installed Product: EoL Date IP': 'EoL Date IP'}, inplace=True)
    if 'Installed Product: EoGS Date IP' in df.columns: # NEW: Rename this column
        df.rename(columns={'Installed Product: EoGS Date IP': 'EoGS Date IP'}, inplace=True)
    if 'Installed Product: Device Age' in df.columns: # NEW: Rename this column
        df.rename(columns={'Installed Product: Device Age': 'Device Age'}, inplace=True)
    if 'Installed Product: Customer/Device Acceptance Date' in df.columns: # NEW: Rename this column
        df.rename(columns={'Installed Product: Customer/Device Acceptance Date': 'Customs Acceptance Date'}, inplace=True)
    if 'Service/Maintenance Contract: Contract Name/Number' in df.columns: # NEW: Rename this column
        df.rename(columns={'Service/Maintenance Contract: Contract Name/Number': 'Contract Name/Number'}, inplace=True)
    if 'Covered Product: Record Number' in df.columns: # NEW: Rename this column
        df.rename(columns={'Covered Product: Record Number': 'Covered Product Record Number'}, inplace=True)
    if 'Current Term Start Date' in df.columns: # NEW: Rename this column
        df.rename(columns={'Current Term Start Date': 'Contract Start Date'}, inplace=True)
    if 'Current Term End Date' in df.columns: # NEW: Rename this column
        df.rename(columns={'Current Term End Date': 'Contract End Date'}, inplace=True)


    # Create 'Display Product Name' for cards/charts
    if 'Installed Product' in df.columns:
        df['Display Product Name'] = df['Installed Product'].apply(
            lambda x: f"{x.split('/')[0]} {x.split('/')[-1]}" if isinstance(x, str) and len(x.split('/')) >= 3 else x
        )
    else:
        df['Display Product Name'] = 'N/A'

    # Convert date columns to datetime objects
    date_cols = ['Warranty Start Date', 'Warranty End Date', 'EoL Date IP', 'EoGS Date IP', 'End Date', 'Start Date', 'Customs Acceptance Date', 'Contract Start Date', 'Contract End Date']
    for col in date_cols:
        if col in df.columns:
            df[col] = pd.to_datetime(df[col], errors='coerce')

    # Calculate Device Age (if not present)
    if 'Device Age' not in df.columns and 'Customs Acceptance Date' in df.columns:
        df['Device Age'] = (datetime.now() - df['Customs Acceptance Date']).dt.days / 365.25
        df['Device Age'] = df['Device Age'].round(1) # Round to 1 decimal place

    # Categorize Device Age into bins
    if 'Device Age' in df.columns:
        # Determine the upper bound for the last bin dynamically but safely
        max_age_in_data = df['Device Age'].max()
        # Ensure the last bin is strictly greater than 10
        # Use a small epsilon to ensure strict increase for pd.cut
        upper_age_bound = max(10.001, max_age_in_data + 0.001)

        df['Device Age Group'] = pd.cut(df['Device Age'],
                                        bins=[0, 5, 10, upper_age_bound],
                                        labels=['0-5 years', '5-10 years', '>10 years'],
                                        right=False, # 0-5 includes 0, excludes 5
                                        include_lowest=True)
        # Fill NaN for devices without age or where age doesn't fit bins
        df['Device Age Group'] = df['Device Age Group'].cat.add_categories('N/A').fillna('N/A')
    else:
        df['Device Age Group'] = 'N/A' # Default if Device Age column is missing

    # Calculate Warranty Remaining (in days)
    df['Warranty Remaining Days'] = (df['Warranty End Date'] - datetime.now()).dt.days

    # Ensure 'Contract Price' is numeric
    if 'Contract Price' in df.columns:
        df['Contract Price'] = pd.to_numeric(df['Contract Price'], errors='coerce').fillna(0)
    else:
        df['Contract Price'] = 0 # Default to 0 if column is missing

    # Ensure 'Weeks To Renewal' is numeric
    if 'Weeks To Renewal' in df.columns:
        df['Weeks To Renewal'] = pd.to_numeric(df['Weeks To Renewal'], errors='coerce').fillna(9999) # Use a large number for non-expiring
    else:
        # If 'Weeks To Renewal' is not provided, try to calculate from 'Contract End Date'
        if 'Contract End Date' in df.columns:
            df['Weeks To Renewal'] = (df['Contract End Date'] - datetime.now()).dt.days / 7
            df['Weeks To Renewal'] = df['Weeks To Renewal'].apply(lambda x: max(0, x)).round(0) # Ensure non-negative and round
        else:
            df['Weeks To Renewal'] = 9999 # Default to a very high number if no contract end date

    # Derive 'Contract Status' if not directly available or to refine
    if 'Contract Status' not in df.columns:
        df['Contract Status'] = 'Active' # Default
    
    # Refine Contract Status based on 'Weeks To Renewal'
    def derive_contract_status(row):
        if row['Weeks To Renewal'] <= 0:
            return 'Expired'
        elif row['Weeks To Renewal'] <= 12:
            return 'Expiring Soon'
        else:
            return 'Active' # Or original status if it exists and is not 'Expired' or 'Expiring Soon'
    df['Contract Status'] = df.apply(derive_contract_status, axis=1)


    # --- Sidebar Filters ---
    st.sidebar.header("Filters")
    
    # Account Filter
    all_accounts = ['All'] + sorted(df['Account'].unique().tolist())
    selected_accounts = st.sidebar.multiselect("Filter by Account", all_accounts, default='All')
    if 'All' in selected_accounts:
        filtered_df = df.copy()
    else:
        filtered_df = df[df['Account'].isin(selected_accounts)].copy()

    # Contract Status Filter
    all_contract_statuses = ['All'] + sorted(df['Contract Status'].unique().tolist())
    selected_contract_statuses = st.sidebar.multiselect("Filter by Contract Status", all_contract_statuses, default='All')
    if 'All' not in selected_contract_statuses:
        filtered_df = filtered_df[filtered_df['Contract Status'].isin(selected_contract_statuses)]

    # Weeks to Renewal Filter
    max_weeks_to_renewal = int(filtered_df['Weeks To Renewal'].max()) if not filtered_df.empty else 0
    weeks_filter_value = st.sidebar.slider(
        "Contracts expiring in next X weeks (0 for expired)",
        min_value=0,
        max_value=max_weeks_to_renewal + 1, # +1 to allow selection up to max
        value=max_weeks_to_renewal + 1, # Default to show all
        step=1
    )
    if weeks_filter_value > 0:
        filtered_df = filtered_df[filtered_df['Weeks To Renewal'] <= weeks_filter_value]
    elif weeks_filter_value == 0: # Show only expired contracts
        filtered_df = filtered_df[filtered_df['Weeks To Renewal'] <= 0]
    
    if filtered_df.empty:
        st.warning("No data matches the selected filters. Please adjust your selections.")
        df_display = pd.DataFrame() # Empty DataFrame to prevent errors
    else:
        df_display = filtered_df.copy()


    st.markdown("---")


    # --- Dashboard Sections ---
    st.header("Overall Fleet & Contract Summary")
    with st.container(border=True): # Wrap KPIs in a container
        col1, col2, col3, col4 = st.columns(4)

        # KPI: Total Devices Installed
        with col1:
            st.metric("Total Devices Installed", df_display.shape[0])

        # KPI: Contracts Expiring Soon (<8 weeks)
        contracts_expiring_soon_kpi = df_display[df_display['Contract Status'] == 'Expiring Soon'].shape[0]
        with col2:
            st.metric("Contracts Expiring Soon (<12 weeks)", contracts_expiring_soon_kpi)

        # KPI: Expired Contracts
        expired_contracts_kpi = df_display[df_display['Contract Status'] == 'Expired'].shape[0]
        with col3:
            st.metric("Expired Contracts", expired_contracts_kpi)

        # KPI: Total Contract Value
        total_contract_value_kpi = df_display['Contract Price'].sum()
        with col4:
            # Assuming currency is consistent, or you can add a currency symbol from a column
            st.metric("Total Contract Value", f"${total_contract_value_kpi:,.0f}")

    st.markdown("---")

    # --- Visualizations ---

    # 1. Contract Status Distribution (Pie Chart)
    st.header("Contract Status Distribution")
    with st.container(border=True):
        if not df_display.empty:
            contract_status_counts = df_display['Contract Status'].value_counts().reset_index()
            contract_status_counts.columns = ['Status', 'Count']
            fig_contract_status = px.pie(contract_status_counts, values='Count', names='Status',
                                         title='Overall Contract Status Distribution',
                                         color_discrete_sequence=COLOR_SEQUENCE,
                                         template="plotly_white")
            st.plotly_chart(fig_contract_status, use_container_width=True)
        else:
            st.info("No data to display for Contract Status Distribution.")

    st.markdown("---")

    # 2. Device Age Distribution (Histogram with Bins)
    st.header("Device Age Distribution")
    with st.container(border=True):
        if 'Device Age Group' in df_display.columns and not df_display.empty:
            # Ensure the order of categories
            age_group_order = ["0-5 years", "5-10 years", ">10 years", "N/A"]
            df_display['Device Age Group'] = pd.Categorical(df_display['Device Age Group'], categories=age_group_order, ordered=True)
            # st.write(df_display) # Removed for cleaner output
            fig_age = px.histogram(df_display, x='Device Age Group',
                                   title='Distribution of Device Ages',
                                   labels={'Device Age Group': 'Device Age (Years)', 'Installed Product': 'Installed Product'},
                                   color_discrete_sequence=COLOR_SEQUENCE,
                                   template="plotly_white",
                                   hover_data=['Serial Number', 'Installed Product', 'Location', 'Account'])
            fig_age.update_layout(bargap=0.8, showlegend=False)
            fig_age.update_xaxes(categoryorder='array', categoryarray=age_group_order)  # Ensure correct order
            st.plotly_chart(fig_age, use_container_width=True)
        else:
            st.info("No data to display for Device Age Distribution or 'Device Age Group' column is missing.")

    st.markdown("---")

    # 3. Upcoming Renewals (Bar Chart)
    st.header("Upcoming Renewals by IP")
    with st.container(border=True):
        if 'Weeks To Renewal' in df_display.columns and not df_display.empty:
            # Filter for contracts expiring in the next 52 weeks (excluding expired ones)
            upcoming_renewals = df_display[df_display['Weeks To Renewal'] <= 52].copy()
            
            if not upcoming_renewals.empty:
                # Create bins for weeks to renewal
                upcoming_renewals['Renewal Period'] = pd.cut(upcoming_renewals['Weeks To Renewal'],
                                                             bins=[-0.1, 0, 12, 26, 52], # Adjusted to capture 0 correctly
                                                             labels=['Expired', '0-12 Weeks', '12-26 Weeks', '26-52 Weeks'],
                                                             right=True,
                                                             include_lowest=True)
                # Add a category for values outside these bins (e.g., >52 weeks if not filtered)
                upcoming_renewals['Renewal Period'] = upcoming_renewals['Renewal Period'].cat.add_categories('>52 Weeks').fillna('>52 Weeks')
                upcoming_renewals = upcoming_renewals[upcoming_renewals['Renewal Period'] != '>52 Weeks'] # Filter out if not needed for chart
                
                renewal_counts = upcoming_renewals.groupby(['Renewal Period','Installed Product' ]).size().unstack(fill_value=0)
                
                fig_renewals = px.bar(renewal_counts, x=renewal_counts.index, y=renewal_counts.columns,
                                      title='Number of Contracts by Upcoming Renewal Period',
                                      labels={'value': 'Number of Contracts', 'Installed Product': 'Installed Product'},
                                      color_discrete_sequence=COLOR_SEQUENCE,
                                      template="plotly_white")
                fig_renewals.update_layout(barmode='stack')
                st.plotly_chart(fig_renewals, use_container_width=True)
            else:
                st.info("No upcoming renewals in the next 52 weeks for the filtered data.")
        else:
            st.info("No data to display for Upcoming Renewals or 'Weeks To Renewal' column is missing.")

    st.markdown("---")

    # 4. Financial Impact (Bar Chart)
    st.header("Contract Value by Location")
    with st.container(border=True):
        if 'Contract Price' in df_display.columns and not df_display.empty:
            contract_value_by_location = df_display.groupby('Location')['Contract Price'].sum().reset_index()
            contract_value_by_location = contract_value_by_location.sort_values(by='Contract Price', ascending=False)
            
            fig_financial = px.bar(contract_value_by_location, y='Location', x='Contract Price',
                                   title='Total Contract Value by Location',
                                   labels={'Contract Price': 'Contract Value'},
                                   color='Location',
                                   color_discrete_sequence=COLOR_SEQUENCE,
                                   template="plotly_white")
            fig_financial.update_layout(xaxis_tickprefix="$", xaxis_tickformat=",.0f", showlegend=False) # Format as currency
            st.plotly_chart(fig_financial, use_container_width=True)
        else:
            st.info("No data to display for Financial Impact or 'Contract Price' column is missing.")

    st.markdown("---")

    # 5. Devices at Risk Table
    st.header("Devices At Risk (Expired or Expiring Soon)")
    with st.container(border=True):
        at_risk_devices = df_display[(df_display['Contract Status'] == 'Expired') | (df_display['Contract Status'] == 'Expiring Soon')].copy()
        if not at_risk_devices.empty:
            # Select relevant columns for the table
            display_cols = ['Account', 'Location', 'Display Product Name', 'Serial Number', 'Contract Status', 'Weeks To Renewal', 'Warranty End Date', 'EoL Date IP']
            table_for_display = at_risk_devices[display_cols].copy()
            
            # Format dates for display
            for col in ['Warranty End Date', 'EoL Date IP']:
                if col in table_for_display.columns:
                    table_for_display[col] = table_for_display[col].dt.strftime('%m/%d/%Y').fillna('N/A')

            st.dataframe(table_for_display, use_container_width=True)
            
            # Download filtered data to Excel
            excel_buffer = BytesIO()
            with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
                df_display.to_excel(writer, sheet_name='Filtered_Data', index=False)
            excel_buffer.seek(0)
            
            st.download_button(
                label="Download Filtered Data as Excel",
                data=excel_buffer,
                file_name="filtered_service_contracts.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        else:
            st.info("No devices currently at risk based on the filters.")

    st.markdown("---")

    # 6. Individual Device Cards (Existing functionality)
    st.header("Machine Fleet Overview")
    cols_per_row = 3
    num_devices = len(df_display)
    num_rows = (num_devices + cols_per_row - 1) // cols_per_row

    for i in range(num_rows):
        cols = st.columns(cols_per_row)
        for j in range(cols_per_row):
            device_idx = i * cols_per_row + j
            if device_idx < num_devices:
                device = df_display.iloc[device_idx]
                with cols[j]:
                    with st.container(border=True):
                        customs_acceptance_date_str = "N/A"
                        if 'Customs Acceptance Date' in device.index and pd.notna(device.get('Customs Acceptance Date')):
                            customs_acceptance_date_str = device['Customs Acceptance Date'].strftime('%m/%d/%Y')

                        original_product_string = device.get('Installed Product', 'Unknown')
                        first_part_for_image = original_product_string.split('/')[0] if isinstance(original_product_string, str) else 'Unknown'
                        sanitized_filename_part = "".join(c if c.isalnum() else '_' for c in first_part_for_image).replace('__', '_').strip('_')
                        image_filename = f"{sanitized_filename_part}.png"

                        base_image_dir = "images/Cards"
                        full_image_path_on_disk = os.path.join(base_image_dir, image_filename)

                        img_col1, img_col2, img_col3 = st.columns([1, 1, 1])
                        with img_col2:
                            if os.path.exists(full_image_path_on_disk):
                                st.image(full_image_path_on_disk, width=100, use_container_width=False, output_format="PNG")
                            else:
                                st.image("https://placehold.co/100x100/A0A0A0/FFFFFF?text=No+Image", width=100, use_container_width=False, output_format="PNG")
                                st.warning(f"Image not found for '{first_part_for_image}'. Looked for: {full_image_path_on_disk}")

                            st.markdown(f"""
                                <div style="text-align: center;">
                                    <h3><i class="fa-solid fa-microchip"></i> {device.get('Display Product Name', 'N/A')}</h3>
                                    <p><i class="fa-solid fa-calendar-xmark"></i> <strong>Contract Expires:</strong> {device.get('End Date', 'N/A').strftime('%m/%d/%Y') if pd.notna(device.get('EoL Date IP')) else 'N/A'}</p>
                                    <p><i class="fa-solid fa-calendar-xmark"></i> <strong>Renewal:</strong> {device.get('Weeks To Renewal', 'N/A')}</p>
                                    <p><i class="fa-solid fa-hourglass-half"></i> <strong>Age:</strong> {device.get('Device Age', 'N/A')} years</p>
                                    <hr></hr>
                                    <p><i class="fa-solid fa-calendar-xmark"></i> <strong>Warranty End Date:</strong> {device.get('Warranty End Date', 'N/A').strftime('%m/%d/%Y') if pd.notna(device.get('Warranty End Date')) else 'N/A'}</p>
                                    <p><i class="fa-solid fa-calendar-xmark"></i> <strong>EoL Date IP:</strong> {device.get('EoL Date IP', 'N/A').strftime('%m/%d/%Y') if pd.notna(device.get('EoL Date IP')) else 'N/A'}</p>
                                    <p><strong>CAT:</strong> {customs_acceptance_date_str}</p>
                                </div>
                            """, unsafe_allow_html=True)
    
    st.sidebar.title('PowerPoint Export')
    if not df_display.empty:
        if st.sidebar.button('Generate Service Contract Slides'):
            with st.spinner('Generating PowerPoint...'):
                try:
                    output_ppt_file = generate_service_contract_slides(df_display, "Radiation Therapy Service Contracts Dashboard")
                    st.write(f"PowerPoint generated: {output_ppt_file}")
                    with open(f"presentations/Service_Agreements/{output_ppt_file}", "rb") as file:
                        st.sidebar.download_button(
                            label="Download Service Contract Slides",
                            data=file,
                            file_name=output_ppt_file,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    st.sidebar.success(f"PowerPoint '{output_ppt_file}' generated successfully!")
                except Exception as e:
                    st.sidebar.error(f"Error generating PowerPoint: {e}")
                    st.sidebar.info("Please ensure 'kaleido' is installed (`pip install kaleido`) for graph export.")
    else:
        st.sidebar.info("Upload data and apply filters to enable PowerPoint generation.")

# This `else` block only executes if `uploaded_file` is None AND `st.session_state['uploaded_df']` is also empty/not set.
else:
    st.markdown("<h3 style='color: rgb(43, 101, 124);'>Load Service Agreement Report</h3>", unsafe_allow_html=True)
    st.write('1. Go to Reports on CLM and select Service Agreement Report (Filter by desired Account).')
    st.markdown("CLM Service Agreement Report: (%s) " % DOWNTIME_URL)
    st.write('2. Select Lightning and Export > Details Only > Format .csv).')
    st.write('3. Upload file.')
    st.write('5. Click generate Power Point.')
